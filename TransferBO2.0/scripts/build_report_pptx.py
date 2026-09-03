"""Build the TransferBO2.0 report PPT from the v2 main-text draft.

Structure follows docs/26_paper_maintext_draft.md (thematic):
  1 cover          2 question & why it matters        3 1.0 lesson / 2.0 design choice
  4 protocol       5 main finding (fig1)              6 BSF curve (fig2)
  7 faster-not-higher (fig3)                          8 what fails: surrogate labels (fig5)
  9 mechanism: rank preservation (fig4)               10 dual channel (fig6)
  11 rejected strategies table                       12 boundaries & abstention
  13 deployment rules                                14 conclusion
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "results" / "presentation"
OUT = ROOT / "results" / "presentation" / "TransferBO2_0_汇报.pptx"

# palette
DARK = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT = RGBColor(0x1F, 0x77, 0xB4)
RED = RGBColor(0xD6, 0x27, 0x28)
GREEN = RGBColor(0x2C, 0xA0, 0x2C)
GRAY = RGBColor(0x59, 0x59, 0x59)
LIGHT = RGBColor(0xF2, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, x, y, w, h, lines, *, size=18, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.15):
    """lines: list of (text, size, bold, color) OR plain strings."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, item in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        if isinstance(item, str):
            item = (item, size, bold, color)
        text, sz, b, c = item
        r = p.add_run()
        r.text = text
        r.font.size = Pt(sz)
        r.font.bold = b
        r.font.color.rgb = c
        r.font.name = "Microsoft YaHei"
    return tb


def add_pic(slide, path, x, y, w=None, h=None):
    kw = {}
    if w: kw["width"] = Inches(w)
    if h: kw["height"] = Inches(h)
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), **kw)


def new_slide(prs, title=None, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, WHITE)
    # header band
    if title:
        band = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.9))  # rectangle
        band.fill.solid(); band.fill.fore_color.rgb = DARK
        band.line.fill.background()
        add_text(slide, 0.45, 0.12, Inches(10.5), Inches(0.7),
                 [(title, 24, True, WHITE)], anchor=MSO_ANCHOR.MIDDLE)
        if subtitle:
            add_text(slide, 0.45, 0.95, Inches(12.0), Inches(0.45),
                     [(subtitle, 13, False, GRAY)])
    return slide


def footer(slide, idx, total):
    add_text(slide, 11.9, 6.55, Inches(1.2), Inches(0.3),
             [(f"{idx} / {total}", 10, False, GRAY)], align=PP_ALIGN.RIGHT)


def main() -> int:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    TOTAL = 14

    # ---------------- 1 cover ----------------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, DARK)
    add_text(s, 0.9, 1.5, Inches(11.5), Inches(1.6), [
        ("Throw most of it away", 44, True, WHITE),
        ("历史 HTE 数据通过五条件清单（而非迁移模型）加速新底物贝叶斯优化", 22, False, RGBColor(0xBD, 0xD7, 0xEE)),
    ])
    add_text(s, 0.9, 3.4, Inches(11.0), Inches(1.4), [
        ("TransferBO 2.0 · 方法学研究汇报", 20, True, ACCENT),
        ("四库 + 一维边界验证 · 71 项任务 · 2,130 条主协议 LOSO 轨迹 · 冻结协议与双对照", 14, False, RGBColor(0xBD, 0xD7, 0xEE)),
        ("2026-08-24 · 证据冻结 / 叙事校正版", 14, False, RGBColor(0xBD, 0xD7, 0xEE)),
    ])
    add_text(s, 0.9, 5.6, Inches(11.5), Inches(1.2), [
        ("核心结论：", 15, True, RGBColor(0xFF, 0xD9, 0x66)),
        ("历史数据最可靠的用法是给出跨底物排序保持的高价值条件清单；", 15, False, WHITE),
        ("历史产率标签不应默认迁入目标 BO 模型；初始化与续跑是两个独立、可选、可弃权的决策。", 15, False, WHITE),
    ])

    # ---------------- 2 question ----------------
    s = new_slide(prs, "研究问题：历史 HTE 数据能否、以及如何加速新底物优化？")
    lines = [
        ("问题 1", 16, True, ACCENT), (" 新底物优化成本高：即使有 HTE，新底物通常需数十次实验才能到好结果", 16, False, DARK),
        ("问题 2", 16, True, ACCENT), (" 历史库充足：同一模板 × 多个底物的稠密条件–产率矩阵随处可见", 16, False, DARK),
        ("问题 3", 16, True, ACCENT), (" 答案并非显然为正：历史测在「别人」身上，底物自身反应活性会移动产率水平", 16, False, DARK),
        ("问题 4", 16, True, ACCENT), (" 文献给出了很多「用法」，但很少在同一冻结协议下、跨多库比较：哪种形式可靠有益？", 16, False, DARK),
    ]
    add_text(s, 0.7, 1.5, Inches(12.0), Inches(2.6), lines, spacing=1.5)
    add_text(s, 0.7, 4.6, Inches(12.0), Inches(1.6), [
        ("我们的定位：不是检验「冷启动 BO 是否强于随机」，而是找一个「正增益的历史数据应用策略」。", 17, True, DARK),
        ("评价口径：主指标 = 优化 AUC（Σ best-so-far，20 步）；双对照 = vs 冷启动 和 vs 随机。", 15, False, GRAY),
    ])

    # ---------------- 3 1.0 lesson ----------------
    s = new_slide(prs, "切入点教训：单源 pair 迁移为负 → 2.0 从设计上改为多源池化")
    add_text(s, 0.7, 1.4, Inches(12.0), Inches(1.2), [
        ("TransferBO 1.0：单源底物对迁移（一个源 → 一个靶），在 Suzuki 模板上效应为负。", 18, True, RED),
        ("原因：单源不仅携带模板通用的条件排序，还携带该底物的特异响应——迁移的「知识」含有噪声。", 15, False, GRAY),
    ])
    add_text(s, 0.7, 2.8, Inches(12.0), Inches(2.8), [
        ("2.0 立项即多源 LOSO（留一底物交叉验证）：对每个靶，历史 = 库内全部其他底物。", 18, True, DARK),
        ("  ", 8, False, DARK),
        ("两个事实（如实陈述）：", 16, True, ACCENT),
        ("① pair 负效应是 1.0 的遗留结果——2.0 没有重跑 pair（pair 轨配置但默认不跑）；", 15, False, DARK),
        ("② 本工作回答的不是「历史有没有用」，而是「多源历史以什么形式进入回路」：", 15, False, DARK),
        ("    作为代理模型标签？作为先验（条件清单）？还是完全不用？", 15, False, DARK),
    ], spacing=1.35)

    # ---------------- 4 protocol ----------------
    s = new_slide(prs, "冻结协议（五库完全一致，预注册后不改）")
    rows = [
        ("评价单位", "leave-one-substrate-out：目标底物外全部底物 = 历史池（多源）"),
        ("初始点", "n_init = 5（前 5 步 = 清单 init；清单 = 历史跨源产率均值 top-5）"),
        ("优化器", "GP (Matern-2.5 ARD) + EI；target-only（历史不经 GP）"),
        ("预算", "B = 20（5 init + 15 续跑）；seeds 0–4"),
        ("指标", "主指标 AUC@20；诊断 AUC@5 / init_best / final_best / 命中 top-5%"),
        ("统计", "seed 平均 → 靶级 → 配对 bootstrap 95% CI（B=5000）；双对照"),
    ]
    y = 1.45
    for k, v in rows:
        add_text(s, 0.7, y, Inches(2.5), Inches(0.4), [(k, 15, True, ACCENT)])
        add_text(s, 3.3, y, Inches(9.4), Inches(0.4), [(v, 15, False, DARK)])
        y += 0.62
    add_text(s, 0.7, 5.6, Inches(12.0), Inches(0.9), [
        ("数据：胺化 15×260（Pd C–N）· 硼化 33×46（Ni C–B）· EDBO Suzuki 12×308（Pd C–C）· HiTEA 11×41–48（Pfizer 独立源）· CHAOS 4×720（一维边界）", 13, False, GRAY),
    ])

    # ---------------- 5 main finding ----------------
    s = new_slide(prs, "主发现：池化 top-5 清单是唯一跨库一致的正策略", "Table 1：四库 × 双对照的 AUC@20 差异（配对标靶 bootstrap 95% CI）")
    add_pic(s, FIG / "fig1_main_forest.png", 0.6, 1.5, w=12.1)
    add_text(s, 0.7, 5.7, Inches(12.0), Inches(1.3), [
        ("四个库方向全为正；两个完整网格库（胺化 +160、硼化 +108）CI 排除 0，88–94% 靶为正。", 15, True, DARK),
        ("EDBO Suzuki vs random 较弱（+92，CI 含 0）——因为该模板冷启动 BO 本身不可靠（随机会更强），属于「可部署性受限」而非「无效」。", 13, False, GRAY),
    ])

    # ---------------- 6 BSF ----------------
    s = new_slide(prs, "增益是「更快」，不是「更高」：第 1 轮即拉开差距", "胺化 best-so-far 曲线（靶级均值，20 步）")
    add_pic(s, FIG / "fig2_bsf_amination.png", 0.7, 1.5, w=7.6)
    add_text(s, 8.6, 1.8, Inches(4.2), Inches(4.2), [
        ("清单把「第 1 轮（前 5 步）」的最好结果直接抬到约 62→66 产率", 15, True, DARK),
        ("  ", 8, False, GRAY),
        ("init_best 差异（vs cold）：", 13, True, ACCENT),
        ("胺化 +12.3 · 硼化 +8.6（均排除 0）", 13, False, DARK),
        ("  ", 8, False, GRAY),
        ("final_best 差异（20 步终点，vs cold）：", 13, True, ACCENT),
        ("胺化 +2.2 · 硼化 +0.3（CI 含 0，≈ 拉平）", 13, False, DARK),
        ("  ", 8, False, GRAY),
        ("历史的价值 = 更好的起点；终点基本不变——「更快到达同一目标」。", 15, True, DARK),
    ], spacing=1.3)

    # ---------------- 7 faster not higher (fig3) ----------------
    s = new_slide(prs, "价值位置是库相关的：init 通道 vs 续跑通道", "Table 3 分解（vs cold，AUC@20）；EDBO Suzuki 例外——优势在后段")
    add_pic(s, FIG / "fig3_init_final.png", 0.6, 1.5, w=8.6)
    add_text(s, 9.4, 1.9, Inches(3.5), Inches(3.8), [
        ("init 主导型", 14, True, ACCENT), ("：胺化、硼化——清单即结论，EI 续跑增益弱（C1 含 0）", 13, False, DARK),
        ("  ", 8, False, GRAY),
        ("后段主导型", 14, True, RED), ("：EDBO Suzuki——清单起点弱（init CI 含 0），但 EI 续跑吃下交互结构（final +5.3 排除 0）", 13, False, DARK),
        ("  ", 8, False, GRAY),
        ("两通道皆弱", 14, True, GRAY), ("：HiTEA 小空间+高噪声，总效应 +26 方向正但 CI 含 0", 13, False, DARK),
    ], spacing=1.3)

    # ---------------- 8 what fails ----------------
    s = new_slide(prs, "什么无效：把历史产率迁入代理模型（null 或有害）")
    add_pic(s, FIG / "fig5_four_arms.png", 0.6, 1.5, w=8.9)
    add_text(s, 9.7, 1.8, Inches(3.3), Inches(4.4), [
        ("sim_weighted（相似度加权）", 14, True, DARK), ("：胺化 +19.3，CI 含 0——null", 13, False, DARK),
        ("  ", 8, False, GRAY),
        ("safe_gate（Spearman 门）", 14, True, DARK), ("：胺化 +11.5，CI 含 0——null", 13, False, DARK),
        ("  ", 8, False, GRAY),
        ("warm 续跑（四臂实验，n=23）", 14, True, RED), ("：历史 warm 点不占预算，却显著变差（B vs A −59.1）", 13, False, DARK),
        ("  ", 8, False, GRAY),
        ("匹配 init 审计（胺化）", 14, True, ACCENT), ("：给定 top-5 起点后 EI 边际仅 +26（含 0）——历史管起点，优化器管精修", 13, False, DARK),
    ], spacing=1.3)

    # ---------------- 9 mechanism ----------------
    s = new_slide(prs, "机制：多源池化 = 排序聚合；排序可迁移，数值不可迁移")
    add_pic(s, FIG / "fig4_rank_pres.png", 0.6, 1.6, w=6.9)
    add_text(s, 7.9, 1.7, Inches(5.1), Inches(4.6), [
        ("化学根源", 16, True, ACCENT),
        ("条件好坏排序 ← 配体（位阻/供电子性）与碱（碱性/溶解性）——对模板内所有底物一致", 14, False, DARK),
        ("产率水平 ← 底物自身反应活性（芳卤电子效应/位阻）——底物特异", 14, False, DARK),
        ("  ", 8, False, GRAY),
        ("因此：排序跨底物保持（ρ 全部为正），数值不可比；极端位阻冲突使保持「部分」化 → 只取顶部 5 个", 14, True, DARK),
        ("  ", 8, False, GRAY),
        ("池化的意义：跨底物排序投票 = 平均掉特异噪声，留下模板通用主效应——这就是「分离通用性质与底物特异性质」的实现", 14, True, DARK),
        ("  ", 8, False, GRAY),
        ("额外证据：CHAOS 一维边界 0.694（五库最高），4/4 靶正——机制不依赖多维条件结构", 13, False, GRAY),
    ], spacing=1.25)

    # ---------------- 10 dual channel ----------------
    s = new_slide(prs, "双通道机制：初始化价值 × 续跑价值 = 独立的两个决策")
    add_pic(s, FIG / "fig6_quadrant.png", 0.6, 1.5, w=6.6)
    # 2x2 deployment table
    cells = [
        ("初始化价值", "续跑价值", "部署建议"),
        ("高", "高", "清单 init + target-only BO（两库都做）"),
        ("高", "低", "只做清单一轮，少续跑/不续跑"),
        ("低", "高", "冷启动/多样化 init + BO"),
        ("低", "低", "弃权：不用这份历史，重新建模或扩大设计空间"),
    ]
    x0, y0 = 7.5, 1.6
    cw, ch = [1.6, 1.6, 3.5], 0.62
    for r, row in enumerate(cells):
        for c, txt in enumerate(row):
            x = x0 + sum(cw[:c]); y = y0 + r * ch
            shp = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(cw[c]), Inches(ch))
            shp.fill.solid()
            shp.fill.fore_color.rgb = LIGHT if r == 0 else WHITE
            shp.line.color.rgb = DARK
            shp.line.width = Pt(0.75)
            tf = shp.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            run = p.add_run(); run.text = txt
            run.font.size = Pt(12 if r == 0 else 11)
            run.font.bold = (r == 0)
            run.font.color.rgb = DARK if r == 0 else (RED if txt == "弃权" else DARK)
            run.font.name = "Microsoft YaHei"
    add_text(s, 7.5, 4.4, Inches(5.4), Inches(1.8), [
        ("事前预测续跑价值（additive R² 分档）已被证伪——暂无可靠事前规则；", 13, False, GRAY),
        ("续跑决策按库型选：init 型库 EI 可选，Suzuki 类 EI 必选；探针测量排序保持是下一步。", 13, False, GRAY),
    ], spacing=1.3)

    # ---------------- 11 rejected strategies ----------------
    s = new_slide(prs, "被 AUC@20 否决的策略（负证据与正证据同框）", "Table 3：每个「自然的直觉」都让数据给出了答案")
    rows = [
        ("相似度加权 pooled GP", "相似底物共享产率水平", "胺化 null（+19.3，CI 含 0）", "不默认使用"),
        ("Spearman 安全门", "少量目标响应可信", "胺化 null（+11.5，CI 含 0）", "当前版不可部署"),
        ("rank 中位数清单", "聚合抗尺度变化", "pooled +1.5（CI 含 0），仅 Suzuki 类边缘", "默认保持 mean"),
        ("历史 warm 进续跑 GP", "更多历史 → 更好后验", "显著变差（B vs A −59.1）", "历史只用于 init"),
        ("additive-R² 续跑规则", "表面结构预测续跑价值", "分档失败（p=0.69）", "无可靠事前规则"),
        ("元特征预测迁移增益", "任务属性可判别增益", "跨库判别 AUC ≈ 0.47（随机）", "不能自动选策略"),
        ("最近邻单源迁移", "最相似底物是最好的供体", "从未超过池化", "池化，不要挑单个"),
    ]
    x0, y0 = 0.6, 1.55
    cw = [3.0, 3.3, 3.4, 2.4]
    ch = 0.62
    heads = ["策略/假设", "最初动机", "AUC@20 结果", "结论"]
    for c, h in enumerate(heads):
        shp = s.shapes.add_shape(1, Inches(x0 + sum(cw[:c])), Inches(y0), Inches(cw[c]), Inches(ch))
        shp.fill.solid(); shp.fill.fore_color.rgb = DARK
        shp.line.fill.background()
        tf = shp.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = h; r.font.size = Pt(12); r.font.bold = True
        r.font.color.rgb = WHITE; r.font.name = "Microsoft YaHei"
    for r_, row in enumerate(rows):
        for c, txt in enumerate(row):
            shp = s.shapes.add_shape(1, Inches(x0 + sum(cw[:c])), Inches(y0 + (r_ + 1) * ch),
                                     Inches(cw[c]), Inches(ch))
            shp.fill.solid(); shp.fill.fore_color.rgb = WHITE if r_ % 2 == 0 else LIGHT
            shp.line.color.rgb = RGBColor(0xBB, 0xBB, 0xBB); shp.line.width = Pt(0.5)
            tf = shp.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            run = p.add_run(); run.text = txt
            run.font.size = Pt(11); run.font.bold = (c == 3)
            run.font.color.rgb = RED if ("否" in txt or "不要" in txt or "弃" in txt) else DARK
            run.font.name = "Microsoft YaHei"

    # ---------------- 12 boundaries ----------------
    s = new_slide(prs, "边界与局限（如实收窄）")
    items = [
        ("跨反应类与跨源", "四库 = 3 个反应类 + 2 个独立数据源（Doyle / Pfizer）；方向一致，但 n=4 库级统计力有限", "强"),
        ("回顾性，非前瞻", "全部为 LOSO 回放；湿实验前瞻验证已预注册（SNAr 模板，128 条件空间）", "中"),
        ("排序保持是相关不是因果", "库级机制支持 + 一维边界一致；尚无靶级可靠预测器", "中"),
        ("未做 plate/batch 校正", "plate_id 是逻辑任务标签；HiTEA 仅一个跨批次样本 → 跨板迁移列为未来工作", "弱"),
        ("统计边界", "5 seeds、单次测量、无噪声重放；靶级 CI 全程报告", "中"),
    ]
    y = 1.5
    for k, v, lvl in items:
        add_text(s, 0.7, y, Inches(2.8), Inches(0.9), [(k, 15, True, DARK)])
        add_text(s, 3.6, y, Inches(8.0), Inches(0.9), [(v, 14, False, GRAY)])
        col = {"强": GREEN, "中": ACCENT, "弱": RED}[lvl]
        add_text(s, 11.7, y, Inches(1.3), Inches(0.5), [(f"证据: {lvl}", 13, True, col)])
        y += 1.0

    # ---------------- 13 deployment ----------------
    s = new_slide(prs, "部署规则（可执行默认）")
    rules = [
        ("1", "源数门槛", "历史底物 ≥3 才启用池化；≥5 推荐；n=1 单源清单不稳定（Jaccard ≈ 0.17）"),
        ("2", "清单", "跨源产率均值排序取 top-5（默认 mean 规则）"),
        ("3", "报告", "逐条件报 source coverage（清单跨源稳定性 0.11–0.40，必须透明）"),
        ("4", "续跑", "按库型：init 型库 EI 可选（清单够用）；Suzuki 类 EI 必选（后段是价值所在）"),
        ("5", "弃权", "两通道都不为正时，不迁移（abstain 是合法默认）"),
        ("6", "口径", "相对指标：AUC@k / 命中 top-5% / 轮次；禁止承诺绝对产率"),
    ]
    y = 1.5
    for num, k, v in rules:
        shp = s.shapes.add_shape(9, Inches(0.7), Inches(y), Inches(0.5), Inches(0.5))  # oval
        shp.fill.solid(); shp.fill.fore_color.rgb = ACCENT; shp.line.fill.background()
        tf = shp.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = num; r.font.size = Pt(15); r.font.bold = True
        r.font.color.rgb = WHITE; r.font.name = "Microsoft YaHei"
        add_text(s, 1.4, y - 0.04, Inches(2.2), Inches(0.5), [(k, 16, True, DARK)])
        add_text(s, 3.6, y - 0.02, Inches(9.2), Inches(0.6), [(v, 14, False, GRAY)])
        y += 0.85

    # ---------------- 14 conclusion ----------------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, DARK)
    add_text(s, 0.9, 0.7, Inches(11.5), Inches(0.8), [
        ("结论：扔掉大部分历史，留下五个条件", 32, True, WHITE),
    ])
    add_text(s, 0.9, 1.8, Inches(11.5), Inches(2.6), [
        ("1. 唯一可靠的正增益做法：多源池化 top-5 条件清单作第 1 轮 + target-only EI（+160 / +108 AUC vs cold，CI 排除 0）", 17, False, WHITE),
        ("2. 历史产率迁入代理模型：null 或负；warm 续跑显著负——初始化与「喂标签」是两个不同的干预", 17, False, WHITE),
        ("3. 化学根基：同一模板内排序由配体/碱的本征性质决定（通用），水平由底物活性决定（特异）→ 排序可迁移、数值不可", 17, False, WHITE),
        ("4. 应用：≥3/≥5 源池化、报 coverage、按库型选续跑、两通道皆弱时弃权", 17, False, WHITE),
    ], spacing=1.5)
    add_text(s, 0.9, 5.3, Inches(11.5), Inches(1.4), [
        ("一句定位：TransferBO 2.0 的贡献不是更复杂的 surrogate 迁移模型，", 15, True, RGBColor(0xFF, 0xD9, 0x66)),
        ("而是在多库序贯优化中证明：历史反应数据最稳健的用途是形成跨底物排序保持的高价值条件清单。", 15, True, RGBColor(0xFF, 0xD9, 0x66)),
        ("下一步：湿实验前瞻验证（已预注册）· SI 表格 · 正文终稿", 14, False, RGBColor(0xBD, 0xD7, 0xEE)),
    ], spacing=1.4)

    prs.save(str(OUT))
    print("saved", OUT)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
