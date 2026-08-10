#!/usr/bin/env python
"""Advisor PPT — bright, polished, design-first (no strict color cap)."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
FIGS = ROOT / "docs" / "figs"
OUT = ROOT / "exports" / "TransferBO_advisor_briefing.pptx"

# Bright design system
INK = RGBColor(0x0F, 0x17, 0x2A)       # deep slate-navy for text
MUTED = RGBColor(0x5B, 0x6B, 0x82)     # secondary text
PAPER = RGBColor(0xFB, 0xFD, 0xFF)     # off-white bg
CARD = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xE6, 0xED, 0xF5)
TEAL = RGBColor(0x10, 0xB8, 0xA6)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
ROSE = RGBColor(0xF4, 0x47, 0x5E)
TEAL_T = RGBColor(0xE3, 0xFA, 0xF4)    # tint
BLUE_T = RGBColor(0xE8, 0xF1, 0xFE)
PURPLE_T = RGBColor(0xF0, 0xE9, 0xFE)
AMBER_T = RGBColor(0xFE, 0xF3, 0xDF)
ROSE_T = RGBColor(0xFE, 0xEC, 0xEF)
FONT = "Microsoft YaHei"


def _ea(run):
    try:
        run._r.rPr.rFonts.set(
            "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}eastAsia", FONT
        )
    except Exception:
        pass


def run_style(run, size=14, bold=False, color=INK, italic=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    _ea(run)


def add_text(slide, left, top, width, height, lines, *, size=14, color=INK,
             bold_first=False, align=PP_ALIGN.LEFT, space_after=3, valign=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    if valign is not None:
        tf.vertical_anchor = valign
    if isinstance(lines, str):
        lines = [lines]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = line
        run_style(run, size=size, bold=(bold_first and i == 0), color=color)
    return box


def add_rect(slide, left, top, width, height, *, fill=CARD, line=LINE, radius=True, shadow=False, line_w=1.0):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    if radius:
        try:
            shp.adjustments[0] = 0.07
        except Exception:
            pass
    return shp


def add_bg(slide, prs):
    add_rect(slide, 0, 0, 13.333, 7.5, fill=PAPER, line=None, radius=False)


def add_footer(slide, page, total):
    add_text(
        slide, 0.5, 7.12, 12.3, 0.28, f"TransferBO · 同库跨任务 BO 实证 · {page:02d} / {total:02d}",
        size=9.5, color=MUTED, align=PP_ALIGN.RIGHT,
    )


def add_pic(slide, path: Path, left, top, width=None, height=None):
    if not path.exists():
        return None
    kw = {}
    if width is not None:
        kw["width"] = Inches(width)
    if height is not None:
        kw["height"] = Inches(height)
    return slide.shapes.add_picture(str(path), Inches(left), Inches(top), **kw)


def header(slide, num: str, kicker: str, title: str, *, accent=TEAL):
    """Eyebrow + big title + accent rule — consistent page header."""
    add_rect(slide, 0.55, 0.42, 0.09, 0.62, fill=accent, line=None)
    add_text(slide, 0.78, 0.32, 3.0, 0.3, kicker, size=11.5, color=MUTED, bold_first=True)
    add_text(slide, 0.78, 0.62, 11.9, 0.55, title, size=22, color=INK, bold_first=True)
    chip = add_rect(slide, 11.9, 0.42, 0.9, 0.4, fill=CARD, line=LINE)
    add_text(slide, 11.9, 0.5, 0.9, 0.3, num, size=11, color=MUTED, align=PP_ALIGN.CENTER)


def card(slide, left, top, width, height, title, lines, *, accent=TEAL, tint=None, size=11.5, title_size=12.5):
    add_rect(slide, left, top, width, height, fill=CARD, line=LINE)
    if tint is not None:
        add_rect(slide, left, top, 0.09, height, fill=accent, line=None, radius=False)
    else:
        add_rect(slide, left, top, width, 0.09, fill=accent, line=None, radius=False)
    add_text(slide, left + 0.22, top + 0.16, width - 0.4, 0.32, title,
             size=title_size, color=accent, bold_first=True)
    if lines:
        add_text(slide, left + 0.22, top + 0.55, width - 0.4, height - 0.65, lines,
                 size=size, color=INK)


def step_chip(slide, left, top, width, title, body, *, accent, tint):
    add_rect(slide, left, top, width, 1.9, fill=tint, line=None)
    add_text(slide, left + 0.16, top + 0.15, width - 0.3, 0.4, title,
             size=14, color=accent, bold_first=True)
    add_text(slide, left + 0.16, top + 0.6, width - 0.3, 1.2, body, size=11, color=INK)


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slides = []

    def new():
        s = prs.slides.add_slide(blank)
        add_bg(s, prs)
        slides.append(s)
        return s

    # ════════ 1 标题页 ════════
    s = new()
    # soft gradient-ish band
    add_rect(s, 0, 0, 13.333, 2.35, fill=BLUE_T, line=None, radius=False)
    add_rect(s, 0, 2.35, 13.333, 0.14, fill=TEAL, line=None, radius=False)
    # floating chips
    for i, (col, tint) in enumerate([(TEAL, TEAL_T), (BLUE, BLUE_T), (PURPLE, PURPLE_T), (AMBER, AMBER_T)]):
        add_rect(s, 10.6 + i * 0.55, 0.55 + (i % 2) * 0.35, 0.34, 0.34, fill=col, line=None)

    add_text(s, 0.75, 0.7, 3.0, 0.35, "TransferBO · 同库跨任务 BO 实证", size=12, color=MUTED, bold_first=True)
    add_text(
        s, 0.75, 1.15, 11.8, 1.0,
        ["历史标签有用吗？", "多样性初始化能否当默认？"],
        size=30, color=INK, bold_first=True,
    )
    add_text(
        s, 0.75, 2.85, 11.8, 0.5,
        "Source-label pooling vs diversity initialisation · 固定 GP–EI 协议下的同库跨任务评估",
        size=15, color=MUTED,
    )
    add_text(
        s, 0.75, 3.5, 11.8, 0.6,
        "定位：严谨的同库跨任务 BO 实证（异质性 / 负结果诚实）\n不是新型 multi-task GP，也不是可部署的“何时迁移”预测器",
        size=13.5, color=INK,
    )
    # roadmap strip
    steps = [
        ("现场", "同库旧标签", TEAL, TEAL_T),
        ("设问", "标签 vs init", PURPLE, PURPLE_T),
        ("设计", "预算对齐", BLUE, BLUE_T),
        ("纠偏", "div ≠ 迁移", AMBER, AMBER_T),
        ("检验", "异质 / 稳健", BLUE, BLUE_T),
        ("收束", "边界清晰", ROSE, ROSE_T),
    ]
    for i, (t, b, col, tint) in enumerate(steps):
        x = 0.75 + i * 2.05
        step_chip(s, x, 4.7, 1.9, t, b, accent=col, tint=tint)
    add_text(s, 0.75, 6.85, 11.8, 0.35, "现场 → 设问 → 设计 → 纠偏 → 检验 → 收束", size=12, color=MUTED)

    # ════════ 2 设问 + 设计 ════════
    s = new()
    header(s, "01", "现场 → 设问 → 设计", "把一句口号拆成两个可检验问题")
    card(
        s, 0.4, 1.35, 4.25, 2.6, "现场混淆（必须拆开）",
        [
            "同库多板：旧板全标 = sunk cost",
            "常混谈三件事：",
            "  (a) 源标签进 GP",
            "  (b) 靶标 diversity init",
            "  (c) 完整 multi-task GP",
        ],
        accent=BLUE, size=11.5,
    )
    card(
        s, 4.8, 1.35, 4.25, 2.6, "两个问题（不是一个口号）",
        [
            "Q1 标签：池化 (x, y_s) 后 frac 是否更高？",
            "    弱/负对有多常见？",
            "Q2 初始化：FPS 相对 random cold 是否仍安全默认？",
            "“何时”= 经验刻画，≠ 预测器",
        ],
        accent=PURPLE, size=11.5,
    )
    card(
        s, 9.2, 1.35, 3.7, 2.6, "公平性先于效果",
        [
            "比：同靶标预算下用不用历史",
            "n_init=20 · B=100 · EI",
            "m=150 = GP 封顶，非魔法比",
            "推断单位 = pair（非 seed）",
        ],
        accent=TEAL, size=11.5,
    )
    # three arms
    card(
        s, 0.4, 4.15, 4.25, 2.6, "cold · 基准",
        ["靶标随机 init；无源标签", "回答：不用历史时的底线"],
        accent=BLUE, size=12,
    )
    card(
        s, 4.8, 4.15, 4.25, 2.6, "diversity · 初始化假说",
        [
            "靶标 FPS；同库 X_s=X_t",
            "≡ cold_diversity（已验证）",
            "负 Δ = init 效应 ≠ 负迁移",
        ],
        accent=AMBER, size=12,
    )
    card(
        s, 9.2, 4.15, 3.7, 2.6, "label · 标签假说",
        [
            "≤150 源标签全程进 GP",
            "无 task ID；联合 z-score",
            "负迁移仅谈这一臂",
        ],
        accent=PURPLE, size=12,
    )

    # ════════ 3 证据① ════════
    s = new()
    header(s, "02", "第一层证据", "设定示意 + 热图：迫使改叙事")
    add_pic(s, FIGS / "fig1_same_library_transfer_v2.png", 0.35, 1.3, width=5.6)
    add_pic(s, FIGS / "fig3_chaos_heatmaps_morgan.png", 6.05, 1.3, width=7.0)
    card(
        s, 0.35, 5.75, 12.6, 1.25, "读图 · 思维转折",
        [
            "左 label 相对 random cold 多为正；右 diversity 多为负。同库下右图 = 靶标 FPS，不是跨板结构迁移 → 负迁移 ⊂ label only。",
            "Morgan / CHAOS 6 对：label Δ≈+0.13 · diversity≈−0.18 · label−diversity≈+0.31（6/6 正）",
        ],
        accent=AMBER, size=12, title_size=12,
    )

    # ════════ 4 证据② ════════
    s = new()
    header(s, "03", "核心证据", "标签通道压过 diversity + 对级异质")
    add_pic(s, FIGS / "fig_label_vs_cold_diversity.png", 0.3, 1.25, width=4.35)
    add_pic(s, FIGS / "fig_si_source_fraction.png", 4.65, 1.25, width=4.55)
    add_pic(s, FIGS / "fig4_chaos_pair_forest_morgan.png", 9.2, 1.25, width=3.95)
    card(
        s, 0.3, 5.65, 12.7, 1.3, "读法（三图一条链）",
        [
            "左：六对 label−diversity 均约 +0.30 → 同库上“用旧标签”优于“默认 diversity”。",
            "中：正迁移对上增加 source_fraction 通常更好。右：pair 森林——强正与弱/近零并存；seed ≠ 推断单位。",
        ],
        accent=TEAL, size=12, title_size=12,
    )

    # ════════ 5 稳健性 ════════
    s = new()
    header(s, "04", "稳健性与边界", "机制 / 描述符 / 外部 / Gate")
    add_pic(s, FIGS / "fig_esi_landscape_scatters.png", 0.3, 1.2, width=5.7)
    add_pic(s, FIGS / "fig5_doyle_pairs_and_targets.png", 6.05, 1.2, width=4.0)
    add_pic(s, FIGS / "fig6_heldout_gate.png", 10.1, 1.2, width=2.9)
    card(
        s, 0.3, 4.9, 4.2, 2.1, "机制：排除错误归因",
        [
            "六对 ρ≈0.60–0.84 均正",
            "Δfrac 与 ρ 几乎无关（≈−0.10）",
            "排除「负迁移=反相关景观」",
        ],
        accent=BLUE, size=11,
    )
    card(
        s, 4.6, 4.9, 4.2, 2.1, "描述符 + SI",
        [
            "CHAOS 三主表示：Morgan / DRFP / xTB",
            "xTB：label +0.15（6/6）· div −0.08",
            "三套符号/排序一致（同级比证）",
        ],
        accent=PURPLE, size=11,
    )
    card(
        s, 8.9, 4.9, 4.1, 2.1, "Doyle + Gate",
        [
            "Doyle 56 对：label≈+0.06 · div≈−0.06",
            "Gate held-out ≈ always-label → No-Go",
            "保留 sanity check，非方法卖点",
        ],
        accent=ROSE, size=11,
    )

    # ════════ 6 收束 ════════
    s = new()
    header(s, "05", "收束", "可以说 / 不可以说 / 下一步")
    claims = [
        (TEAL, TEAL_T, "可以说", "同库同预算下，源标签池化平均优于随机冷启动；存在弱/负对。"),
        (TEAL, TEAL_T, "可以说", "本协议 FPS diversity 平均不如 random cold；同库下不是跨板结构迁移。"),
        (PURPLE, PURPLE_T, "必须说", "推断看 pair；负迁移仅 label；label ≠ 完整 MTGP；联合 z-score 有局限。"),
        (ROSE, ROSE_T, "不可以说", "可部署 when-to-transfer；diversity 永远有害；机制已被景观相关解释。"),
    ]
    y = 1.35
    for col, tint, tag, text in claims:
        add_rect(s, 0.4, y, 12.5, 0.82, fill=tint, line=None)
        add_rect(s, 0.4, y, 0.09, 0.82, fill=col, line=None, radius=False)
        add_text(s, 0.65, y + 0.2, 1.6, 0.4, tag, size=13.5, color=col, bold_first=True)
        add_text(s, 2.4, y + 0.2, 10.2, 0.5, text, size=13, color=INK)
        y += 0.94
    card(
        s, 0.4, 5.25, 6.15, 1.9, "本篇该停",
        [
            "叙事/术语已纠偏（稿 v0.4b）",
            "CHAOS + Doyle + xTB + SI 齐",
            "禁止再扩四板同质大网格",
        ],
        accent=BLUE, size=12,
    )
    card(
        s, 6.75, 5.25, 6.15, 1.9, "下一篇 / 不挡投",
        [
            "真 task-aware MTGP",
            "共享靶标 init 公平性补跑（可选）",
            "机制需要更多任务，非更多同质对",
        ],
        accent=PURPLE, size=12,
    )

    # ════════ 7 备忘 ════════
    s = new()
    header(s, "06", "一页备忘", "思维进程 × 数字锚点")
    card(
        s, 0.4, 1.3, 6.3, 5.85, "思维进程（口述可用）",
        [
            "1 现场：同库旧标签要不要用？diversity 能否默认？",
            "2 拆问：标签通道 vs 初始化通道；不做完整 MTGP。",
            "3 设计：靶标预算对齐；pair 推断；协议锁定。",
            "4 纠偏：diversity ≡ cold_diversity；NTR 仅 label。",
            "5 检验：异质森林 → 景观假说排除 → xTB / Doyle",
            "   → Gate No-Go（学不会何时）。",
            "6 收束：平均有用 ≠ 总是迁；边界写清再投稿。",
            "",
            "材料：manuscript_draft_DD.md · protocol.yaml",
            "      results/paper_stats/ · data/descriptors/",
        ],
        accent=BLUE, size=13,
    )
    card(
        s, 6.9, 1.3, 6.0, 5.85, "数字锚点（汇报用）",
        [
            "CHAOS Morgan（6 对）",
            "  label vs cold ≈ +0.13",
            "  diversity vs cold ≈ −0.18",
            "  label − diversity ≈ +0.31（6/6）",
            "",
            "CHAOS xTB（与 FP 同级）",
            "  label +0.15（6/6） · div −0.08 · label−div +0.23（6/6）",
            "",
            "Doyle 56 对",
            "  label ≈ +0.064 · div ≈ −0.059",
            "",
            "景观：ρ≈0.60–0.84 · corr(Δ,ρ)≈−0.10",
            "Gate：held-out ≈ always-label → No-Go",
        ],
        accent=PURPLE, size=13,
    )

    total = len(slides)
    for i, slide in enumerate(slides, 1):
        add_footer(slide, i, total)

    out = OUT
    out.parent.mkdir(parents=True, exist_ok=True)
    try:
        prs.save(str(out))
        print(f"Wrote {out}")
    except PermissionError:
        fallback = out.with_name(out.stem + "_new.pptx")
        prs.save(str(fallback))
        print(f"canonical locked; wrote {fallback}")
        out = fallback
    print(f"slides={total}")
    return out


if __name__ == "__main__":
    build()
