#!/usr/bin/env python
"""Research-journey briefing PPT — from project start through EDBO pivot."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
FIGS = ROOT / "docs" / "figs"
OUT = ROOT / "exports" / "TransferBO_research_journey_20260805.pptx"

INK = RGBColor(0x14, 0x1C, 0x2B)
MUTED = RGBColor(0x5A, 0x67, 0x7A)
PAPER = RGBColor(0xF7, 0xF5, 0xF0)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xD9, 0xE0, 0xE8)
TEAL = RGBColor(0x0D, 0x94, 0x88)
BLUE = RGBColor(0x2F, 0x6F, 0xED)
AMBER = RGBColor(0xC9, 0x7A, 0x1A)
ROSE = RGBColor(0xC2, 0x3B, 0x4A)
TEAL_T = RGBColor(0xE6, 0xF6, 0xF3)
BLUE_T = RGBColor(0xEA, 0xF1, 0xFC)
AMBER_T = RGBColor(0xFD, 0xF3, 0xE3)
ROSE_T = RGBColor(0xFB, 0xEB, 0xED)
DARK = RGBColor(0x12, 0x2A, 0x3A)
FONT = "Microsoft YaHei"


def _ea(run):
    try:
        run._r.rPr.rFonts.set(
            "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}eastAsia", FONT
        )
    except Exception:
        pass


def run_style(run, size=14, bold=False, color=INK, italic=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    _ea(run)


def add_text(
    slide,
    left,
    top,
    width,
    height,
    lines,
    *,
    size=14,
    color=INK,
    bold_first=False,
    align=PP_ALIGN.LEFT,
    space_after=4,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    if isinstance(lines, str):
        lines = [lines]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = line
        run_style(run, size=size, bold=(bold_first and i == 0), color=color)
    return box


def add_rect(slide, left, top, width, height, *, fill=CARD, line=LINE, radius=True):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1.0)
    if radius:
        try:
            shp.adjustments[0] = 0.06
        except Exception:
            pass
    return shp


def add_bg(slide):
    add_rect(slide, 0, 0, 13.333, 7.5, fill=PAPER, line=None, radius=False)


def add_accent_bar(slide, color=TEAL):
    add_rect(slide, 0, 0, 0.12, 7.5, fill=color, line=None, radius=False)


def add_footer(slide, page, total):
    add_text(
        slide,
        0.5,
        7.15,
        12.3,
        0.28,
        f"TransferBO · 从启动到主线纠偏 · 2026-08-05 · {page:02d}/{total:02d}",
        size=10,
        color=MUTED,
    )


def add_title(slide, title, subtitle=None):
    add_text(slide, 0.5, 0.28, 12.2, 0.45, title, size=22, bold_first=True, color=INK)
    if subtitle:
        add_text(slide, 0.5, 0.70, 12.2, 0.35, subtitle, size=12, color=MUTED)


def add_pic(slide, path: Path, left, top, width, height=None):
    if not path.is_file():
        add_rect(slide, left, top, width, height or 3.5, fill=BLUE_T)
        add_text(
            slide,
            left + 0.2,
            top + 1.2,
            width - 0.4,
            0.8,
            f"[缺图] {path.name}",
            size=12,
            color=MUTED,
            align=PP_ALIGN.CENTER,
        )
        return
    kwargs = {"width": Inches(width)}
    if height is not None:
        kwargs["height"] = Inches(height)
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), **kwargs)


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_bar(slide)
    return slide


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # —— 1 封面 ——
    s = new_slide(prs)
    add_rect(s, 0, 0, 13.333, 7.5, fill=DARK, line=None, radius=False)
    add_text(s, 0.8, 1.6, 11.5, 0.5, "方向三 · TransferBO", size=16, color=TEAL, bold_first=True)
    add_text(
        s,
        0.8,
        2.2,
        11.5,
        1.2,
        "跨反应板 Transfer / Warm-start BO：项目历程与阶段结果",
        size=30,
        bold_first=True,
        color=CARD,
    )
    add_text(
        s,
        0.8,
        3.7,
        11.5,
        2.0,
        [
            "从组内「方向三」立项 → 同库四板迁移网格 → TransferGate → 外部验证",
            "→ 发现一维局限 → EDBO 多维条件库主线纠偏",
            "",
            "汇报日期：2026-08-05",
        ],
        size=15,
        color=RGBColor(0xD0, 0xDC, 0xE6),
        space_after=6,
    )

    # —— 2 总览时间线（全程）——
    s = new_slide(prs)
    add_title(s, "全程时间线：从启动到现在", "六段主弧，后面分页展开")
    phases = [
        ("① 立项", "方向三定位\n纯计算 · 跨板迁移\n反哺湿实验冷启动", TEAL_T, TEAL),
        ("② 基建", "W1：数据/协议/BO\nCHAOS 四板入库\n冷启动冒烟", BLUE_T, BLUE),
        ("③ 主网格", "W3：策略×表示\ncold / diversity / label\npair 级统计", AMBER_T, AMBER),
        ("④ Gate", "W6–W8 TransferGate\nheld-out plate_4\n结论：No-Go", ROSE_T, ROSE),
        ("⑤ 外部", "Doyle2018 验证\nSURF 审计拒绝\n效应更小更异质", TEAL_T, TEAL),
        ("⑥ 纠偏", "PK=一维补充\nEDBO=主线\nhelp/null/harm", BLUE_T, BLUE),
    ]
    for i, (t, b, fill, acc) in enumerate(phases):
        x = 0.35 + i * 2.15
        add_rect(s, x, 1.45, 2.05, 5.1, fill=fill)
        add_rect(s, x, 1.45, 2.05, 0.12, fill=acc, line=None, radius=False)
        add_text(s, x + 0.12, 1.8, 1.8, 0.7, t, size=16, bold_first=True, color=INK)
        add_text(s, x + 0.12, 2.7, 1.8, 3.4, b, size=12, color=MUTED, space_after=5)

    # —— 3 立项动机 ——
    s = new_slide(prs)
    add_title(s, "① 项目启动：为什么做 TransferBO", "组内「方向三」——不占机械臂的纯计算跨板迁移研究")
    add_rect(s, 0.5, 1.3, 6.1, 5.3, fill=CARD)
    add_text(
        s,
        0.75,
        1.55,
        5.6,
        4.8,
        [
            "一句话问题（立项时）",
            "反应板 A 的结果，能否帮助板 B",
            "用更少查询找到高响应条件？",
            "迁什么？何时负迁移？",
            "",
            "与另外两条线的分工",
            "• 方向一：湿实验闭环（占机器人）",
            "• 方向二：描述符 × BO（湿/算）",
            "• 方向三：公开 HTE 上的迁移协议",
            "  → 反哺一/二的冷启动",
        ],
        size=14,
        bold_first=True,
        space_after=5,
    )
    add_rect(s, 6.85, 1.3, 5.95, 5.3, fill=TEAL_T)
    add_text(
        s,
        7.1,
        1.55,
        5.5,
        4.8,
        [
            "期刊阶梯（计划书）",
            "• 保底：Digital Discovery 实证",
            "• 抬高：外部域验证 + 机制",
            "• 冲顶：基准包 / 可选湿验证",
            "",
            "早期抬高叙事：TransferGate",
            "优化前预测「该不该迁、迁什么」",
            "——后来 W8 判定 No-Go，",
            "诚实降为非主卖点。",
            "",
            "文档：方向三-TransferBO纯计算方案.md",
        ],
        size=14,
        bold_first=True,
        space_after=5,
    )

    # —— 4 W1 基建 ——
    s = new_slide(prs)
    add_title(s, "② 基建期（W1）：把回顾性 BO 跑通", "公开多板 HTE → 统一表 → GP–EI 离散循环 → 可复现配置")
    boxes = [
        ("数据", "CHAOS/PK 四板添加剂\n720 × 4；UV210\n统一 plate CSV"),
        ("Oracle", "查表 = 做实验\n离散候选池\n预算内逐点揭示"),
        ("代理/采集", "sklearn GP\nMatérn + White\n主采集 EI"),
        ("策略原型", "cold_start\nlabel_warm\ndiversity / multitask"),
    ]
    for i, (t, b) in enumerate(boxes):
        x = 0.45 + i * 3.2
        add_rect(s, x, 1.4, 3.05, 3.4, fill=CARD)
        add_text(s, x + 0.2, 1.7, 2.65, 0.5, t, size=18, bold_first=True, color=TEAL)
        add_text(s, x + 0.2, 2.5, 2.65, 2.0, b, size=13, color=MUTED, space_after=5)
    add_rect(s, 0.45, 5.05, 12.4, 1.5, fill=BLUE_T)
    add_text(
        s,
        0.7,
        5.3,
        11.9,
        1.0,
        [
            "协议萌芽：n_init / budget / ≥20 seeds；held-out plate_4 冻结到 Gate 评测；",
            "实现落点：src/transferbo + configs/protocol.yaml + run_experiment / run_transfer_grid",
        ],
        size=14,
        color=INK,
        space_after=6,
    )

    # —— 5 科学设定示意 ——
    s = new_slide(prs)
    add_title(s, "同库跨任务设定（贯穿全程的操作定义）", "共享离散候选库 X；任务随板/底物更换；比较是否使用历史标签")
    add_pic(s, FIGS / "fig1_same_library_transfer_v2.png", 1.3, 1.2, 10.7, 5.5)

    # —— 6 W3 网格设计 ——
    s = new_slide(prs)
    add_title(s, "③ 主实验设计（W3）：在四板上做迁移消融", "相对 Ranković 单板 BO：我们问「有了另一张全标板之后该怎么迁」")
    add_rect(s, 0.5, 1.3, 4.1, 5.3, fill=CARD)
    add_text(
        s,
        0.75,
        1.55,
        3.6,
        4.8,
        [
            "策略（信息类型）",
            "• cold：目标随机冷启动",
            "• diversity：目标侧 FPS",
            "• label：源标签池化进 GP",
            "• multitask：与 label 等价实现",
            "",
            "表示",
            "• Morgan / DRFP /（后加）xTB",
            "",
            "推断单元",
            "• source→target pair",
            "• seeds 量化优化噪声",
        ],
        size=13,
        bold_first=True,
        space_after=4,
    )
    add_rect(s, 4.85, 1.3, 7.95, 5.3, fill=AMBER_T)
    add_text(
        s,
        5.15,
        1.55,
        7.4,
        4.8,
        [
            "关键冻结",
            "• 开发折：plates 1–3（6 个有向对）",
            "• held-out：plate_4（Gate）",
            "• max_warm_points=150（GP 可算）",
            "• 主指标：frac-of-opt、Δfrac、曲线",
            "",
            "文献定位（当时）",
            "• 不是重复 CHAOS 单板选配置",
            "• 而是同库跨板：迁标签 vs 迁结构",
            "• 强调负迁移与异质性（相对 MTBO 叙事）",
            "",
            "2026-07 末：20 seeds 全网格落地",
        ],
        size=13,
        bold_first=True,
        space_after=4,
    )

    # —— 7 PK/CHAOS 结果图 ——
    s = new_slide(prs)
    add_title(
        s,
        "③ 四板主结果（现称 PK2022）：label 平均正，diversity 平均负",
        "这是项目第一座里程碑——也是后来需要降维解释的结果",
    )
    add_pic(s, FIGS / "fig3_chaos_heatmaps_morgan.png", 0.35, 1.15, 6.4, 5.5)
    add_pic(s, FIGS / "fig4_chaos_pair_forest_morgan.png", 6.9, 1.15, 5.95, 5.5)

    # —— 8 数字摘要 + diversity ——
    s = new_slide(prs)
    add_title(s, "③ 结果要点与多样性对照", "pair 级异质；diversity 不能当跨板默认")
    add_rect(s, 0.5, 1.3, 6.1, 5.3, fill=CARD)
    add_text(
        s,
        0.75,
        1.55,
        5.6,
        4.8,
        [
            "开发折 overall（摘录）",
            "• label Morgan Δfrac ≈ +0.13",
            "• label DRFP  ≈ +0.15",
            "• label xTB   ≈ +0.15",
            "• diversity 多为负（约 −0.15~−0.18）",
            "",
            "解释边界（当时已写）",
            "• 平均正 ≠ 每个 pair 都正",
            "• diversity_warm ≡ 目标侧 FPS",
            "  （同库同编码，不是源特异结构迁移）",
            "• NTR 仅定义在 label vs cold",
        ],
        size=14,
        bold_first=True,
        space_after=5,
    )
    add_pic(s, FIGS / "fig_label_vs_cold_diversity.png", 6.85, 1.3, 5.95, 5.3)

    # —— 9 Gate ——
    s = new_slide(prs)
    add_title(s, "④ TransferGate（W6–W8）：抬高尝试与 No-Go", "想事前预测「该不该迁」——held-out 上未超过 always-label")
    add_rect(s, 0.5, 1.3, 6.1, 5.3, fill=CARD)
    add_text(
        s,
        0.75,
        1.55,
        5.6,
        4.8,
        [
            "动机",
            "把现象报告升级为可部署判据：",
            "未见目标标签时，预测迁移是否值得。",
            "",
            "做法",
            "• 用开发折 pair 特征训 Gate",
            "• plate_4 held-out 终评",
            "• 与 always-label / cold 比",
            "",
            "结论（W8）",
            "• No-Go：不作为主卖点",
            "• 样本极少，易塌成「多数开 label」",
            "• 诚实负结果保留进论文局限",
        ],
        size=13,
        bold_first=True,
        space_after=4,
    )
    add_pic(s, FIGS / "fig6_heldout_gate.png", 6.85, 1.3, 5.95, 5.3)

    # —— 10 Doyle ——
    s = new_slide(prs)
    add_title(
        s,
        "⑤ 外部验证：Doyle2018 Buchwald–Hartwig",
        "同构设定（条件库×换底物）；OHE；抬高档必需的独立域",
    )
    add_pic(s, FIGS / "fig5_doyle_pairs_and_targets.png", 0.35, 1.15, 6.4, 5.5)
    add_pic(s, FIGS / "fig5_doyle_target_bars.png", 6.9, 1.15, 5.95, 5.5)

    # —— 11 Doyle 文字 ——
    s = new_slide(prs)
    add_title(s, "⑤ Doyle 结果如何读", "方向可作确认；效应量与表示不可与指纹/DFT 硬比")
    cards = [
        ("设定", "15 底物中预注册 8 个\n240 共享条件\nligand×base×additive\n表示：条件键 OHE"),
        ("结果", "label 平均 Δfrac 为正但更小\n异质性更强\n部分靶难/天花板"),
        ("审计", "SURF 等未过 same-library\n拒绝硬塞进主文验证"),
        ("当时叙事", "「外部同向」支撑\n四板结论可推广一截\n——随后发现还需再纠"),
    ]
    for i, (t, b) in enumerate(cards):
        x = 0.45 + (i % 2) * 6.4
        y = 1.35 + (i // 2) * 2.7
        add_rect(s, x, y, 6.15, 2.5, fill=CARD)
        add_text(s, x + 0.25, y + 0.3, 5.7, 0.5, t, size=18, bold_first=True, color=BLUE)
        add_text(s, x + 0.25, y + 1.0, 5.7, 1.2, b, size=13, color=MUTED, space_after=4)

    # —— 12 纠偏转折 ——
    s = new_slide(prs)
    add_title(s, "⑥ 关键转折：主战场其实问错了维度", "2026-08：意识到「CHAOS/PK 四板」是一维添加剂空间")
    add_rect(s, 0.5, 1.35, 6.1, 5.2, fill=ROSE_T)
    add_text(
        s,
        0.75,
        1.65,
        5.6,
        4.6,
        [
            "问题在哪",
            "X_PK ≈ {additive}",
            "板间 = 反应变体差异",
            "≠ 多维条件优化 × 换底物",
            "",
            "若继续以 PK 为「主证据」",
            "会把一维筛上的标签复用",
            "误写成条件库迁移通论。",
        ],
        size=15,
        bold_first=True,
        space_after=6,
    )
    add_rect(s, 6.85, 1.35, 5.95, 5.2, fill=TEAL_T)
    add_text(
        s,
        7.1,
        1.65,
        5.5,
        4.6,
        [
            "纠偏后的三库角色",
            "• EDBO2021 = 主线 / 主证据",
            "• Doyle2018 = 外部验证",
            "• PK2022 = 一维补充",
            "",
            "数据溯源",
            "• 库名用 Prieto–Kullmer 2022",
            "• CHAOS=Ranković 2024 复用/BO 文",
            "",
            "题眼改为：help / null / harm",
        ],
        size=15,
        bold_first=True,
        space_after=6,
    )

    # —— 13 EDBO 设计 ——
    s = new_slide(prs)
    add_title(s, "⑥ 新主线：EDBO Suzuki 多维条件库", "任务=底物对；X=L×B×Sol（308）；Morgan / DRFP / DFT")
    add_rect(s, 0.5, 1.3, 5.0, 5.3, fill=CARD)
    add_text(
        s,
        0.75,
        1.55,
        4.5,
        4.8,
        [
            "为什么是主证据",
            "• 真·多维条件组合空间",
            "• 跨底物优化最贴近应用",
            "• 与 Doyle 同构、表示更强",
            "",
            "预注册 8 任务",
            "suz_t12→t8→t1→t9",
            "→t10→t7→t5→t3",
            "",
            "协议",
            "cold vs label only",
            "n_init=20, B=100, 20 seeds",
        ],
        size=14,
        bold_first=True,
        space_after=4,
    )
    add_rect(s, 5.75, 1.3, 7.05, 5.3, fill=BLUE_T)
    add_text(
        s,
        6.05,
        1.55,
        6.5,
        4.8,
        [
            "计算进度（2026-08-05）",
            "",
            "✓ cold × Morgan / DRFP / DFT",
            "✓ label × Morgan（1120/1120）",
            "→ label × DFT 优先收尾（查表更快）",
            "⏸ label × DRFP 待续",
            "",
            "已落盘",
            "• 条件 DFT 表；Morgan/DRFP 指纹表",
            "• PCA-128（未接主协议）",
            "",
            "氨化：暂停，Suzuki 后再议",
        ],
        size=14,
        bold_first=True,
        space_after=5,
    )

    # —— 14 cold 曲线 ——
    s = new_slide(prs)
    add_title(s, "EDBO · Cold 基线学习曲线", "随机 init 20 次后进入 GP–EI；三表示后期接近")
    add_pic(s, FIGS / "fig_edbo_suzuki_cold_learning_curves_yield.png", 1.5, 1.15, 10.2, 5.6)

    # —— 15 Morgan 总曲线 ——
    s = new_slide(prs)
    add_title(
        s,
        "EDBO · Morgan：label vs cold（阶段主结果）",
        "总平均：label 未显示稳定优势——与 PK/Doyle「平均正」叙事冲突，触发问题重写",
    )
    add_pic(s, FIGS / "fig_edbo_suzuki_morgan_label_vs_cold_yield.png", 1.5, 1.15, 10.2, 5.6)

    # —— 16 分板 ——
    s = new_slide(prs)
    add_title(
        s,
        "EDBO · Morgan：按目标板",
        "8 板分面同样近零/略负——不是总平均制造的假象",
    )
    add_pic(s, FIGS / "fig_edbo_suzuki_morgan_label_vs_cold_by_target.png", 0.55, 1.1, 12.2, 5.7)

    # —— 17 解读 ——
    s = new_slide(prs)
    add_title(s, "如何理解「主线变负/近零」", "项目升级，不是推倒")
    points = [
        ("从平均增益 → 条件问题", "问何时 help / null / harm，比坚持「平均有效」更诚实、更好发"),
        ("共享 X ≠ 共享 y(x)", "无 task ID 的 raw pooling 假设过强；源点可主导早期 GP"),
        ("Cold 很强 / 有天花板", "B≈50 已高产率时，transfer 难转化为最终 Δfrac"),
        ("旧库结果仍有位置", "PK=一维补充；Doyle=外部验证；不再当通论支柱"),
    ]
    for i, (t, b) in enumerate(points):
        y = 1.25 + i * 1.35
        add_rect(s, 0.5, y, 12.3, 1.2, fill=CARD)
        add_text(s, 0.8, y + 0.2, 11.7, 0.35, t, size=16, bold_first=True, color=TEAL)
        add_text(s, 0.8, y + 0.6, 11.7, 0.4, b, size=13, color=MUTED)

    # —— 18 方法阶梯 ——
    s = new_slide(prs)
    add_title(
        s,
        "下一步：预注册的最小方法阶梯",
        "避免审稿概括为「只证了一个过强 naive pooling 不工作」",
    )
    arms = [
        ("A0", "Cold", "不用旧标签", TEAL),
        ("A1", "Raw pooling", "直接混 raw yield\n（当前主协议）", ROSE),
        ("A2", "Rank / norm", "任务内相对尺度\n后再池化", AMBER),
        ("A3", "Weighted", "源标签降权\n弱证据使用", BLUE),
    ]
    for i, (c, n, d, col) in enumerate(arms):
        x = 0.45 + i * 3.2
        add_rect(s, x, 1.45, 3.05, 4.3, fill=CARD)
        add_rect(s, x, 1.45, 3.05, 0.12, fill=col, line=None, radius=False)
        add_text(s, x + 0.2, 1.85, 2.65, 0.4, c, size=12, color=col, bold_first=True)
        add_text(s, x + 0.2, 2.4, 2.65, 0.6, n, size=18, bold_first=True)
        add_text(s, x + 0.2, 3.3, 2.65, 1.8, d, size=14, color=MUTED, space_after=5)
    add_text(
        s,
        0.5,
        6.05,
        12.3,
        0.7,
        "必做 S0：matched-target-init · 故事：旧数据如何安全使用 · 禁止为出正结果加臂",
        size=13,
        color=MUTED,
    )

    # —— 19 状态 ——
    s = new_slide(prs)
    add_title(s, "任务状态（启动至今）", "详见 docs/待处理事项.md")
    cols = [
        (
            "已做",
            TEAL_T,
            [
                "立项与代码框架",
                "四板网格 + 三表示",
                "Gate No-Go",
                "Doyle 外部验证",
                "一维纠偏与主张锁定",
                "EDBO 划板与 cold×3",
                "label×Morgan 全量",
                "Morgan 总/分板图",
                "方法阶梯预注册",
            ],
        ),
        (
            "正在做",
            AMBER_T,
            [
                "label×DFT 优先",
                "",
                "紧接着：",
                "续 DRFP",
                "三表示汇总主图",
                "C1 结论落锤",
            ],
        ),
        (
            "待做",
            BLUE_T,
            [
                "S0 matched-init",
                "A2 rank-norm",
                "A3 source-weight",
                "升版稿（不覆盖 v0.4c）",
                "氨化（点头再开）",
                "",
                "降权：ICM / 复杂 MTGP",
            ],
        ),
    ]
    for i, (title, fill, lines) in enumerate(cols):
        x = 0.45 + i * 4.25
        add_rect(s, x, 1.3, 4.05, 5.3, fill=fill)
        add_text(s, x + 0.25, 1.55, 3.55, 0.5, title, size=18, bold_first=True)
        add_text(s, x + 0.25, 2.2, 3.55, 4.0, ["• " + x for x in lines if x] , size=13, space_after=5)

    # —— 20 收束 ——
    s = new_slide(prs)
    add_title(s, "一条连贯故事线", "从启动问题到当前主线，逻辑不断裂")
    add_rect(s, 0.5, 1.35, 12.3, 5.2, fill=CARD)
    add_text(
        s,
        0.85,
        1.7,
        11.6,
        4.6,
        [
            "1. 立项：公开 HTE 上检验跨板迁移——迁什么、何时负迁移；并尝试 Gate。",
            "2. 交付：四板上 label 平均有益、diversity 平均有害；Gate held-out No-Go；Doyle 外部弱确认。",
            "3. 纠偏：四板主战场实为一维添加剂空间；主证据改到 EDBO 多维条件库 × 换底物。",
            "4. 现状：EDBO×Morgan 上 raw pooling 未稳定优于 cold——支持重写题眼，而非否定前序工作。",
            "5. 前方：收齐 DFT/DRFP；用 Cold→Raw→Rank-norm→Weighted 阶梯回答「旧数据如何安全使用」。",
        ],
        size=15,
        space_after=14,
    )

    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        if i == 1:
            continue
        add_footer(slide, i, total)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    print(f"wrote {build()}")
