#!/usr/bin/env python
"""Build advisor briefing PPT + DOCX for locked EDBO Suzuki main line (v0.7).

PPT structure: experimental design → purpose → results → data analysis.
v0.7 adds W8 amination full-grid (family-dependent boundary).
"""

from __future__ import annotations

from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
from docx.oxml.ns import qn
from docx.shared import Cm, Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PRGB
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches as PInches
from pptx.util import Pt as PPt

ROOT = Path(__file__).resolve().parents[1]
FIGS = ROOT / "docs" / "figs" / "main"
ESI = ROOT / "docs" / "figs"
OUT_DIR = ROOT / "docs" / "briefings"
PPT_OUT = OUT_DIR / "汇报_TransferBO_EDBO_Suzuki_v0.7.pptx"
DOC_OUT = OUT_DIR / "汇报_TransferBO_EDBO_Suzuki_v0.7.docx"
PPT_ALIAS = OUT_DIR / "汇报_TransferBO_EDBO_Suzuki_v0.6.pptx"
DOC_ALIAS = OUT_DIR / "汇报_TransferBO_EDBO_Suzuki_v0.6.docx"

# Design — clean academic, not purple-glow AI default
INK = PRGB(0x1A, 0x1F, 0x2E)
MUTED = PRGB(0x5A, 0x64, 0x72)
PAPER = PRGB(0xF7, 0xF5, 0xF0)
CARD = PRGB(0xFF, 0xFF, 0xFF)
LINE = PRGB(0xD8, 0xD2, 0xC8)
ACCENT = PRGB(0x1F, 0x6B, 0x5B)  # deep teal-green
ACCENT2 = PRGB(0xB5, 0x3D, 0x2B)  # terracotta for negative/warn
NAVY = PRGB(0x2F, 0x4B, 0x7C)
FONT = "Microsoft YaHei"
TOTAL = 22


def _ea(run):
    try:
        run._r.rPr.rFonts.set(
            "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}eastAsia", FONT
        )
    except Exception:
        pass


def run_style(run, size=14, bold=False, color=INK, italic=False):
    run.font.name = FONT
    run.font.size = PPt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    _ea(run)


def add_text(
    slide,
    left,
    top,
    width,
    height,
    lines,
    *,
    size=14,
    color=INK,
    bold_first=False,
    align=PP_ALIGN.LEFT,
    space_after=4,
    valign=None,
):
    box = slide.shapes.add_textbox(PInches(left), PInches(top), PInches(width), PInches(height))
    tf = box.text_frame
    tf.word_wrap = True
    if valign is not None:
        tf.vertical_anchor = valign
    if isinstance(lines, str):
        lines = [lines]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = PPt(space_after)
        run = p.add_run()
        run.text = line
        run_style(run, size=size, bold=(bold_first and i == 0), color=color)
    return box


def add_rect(slide, left, top, width, height, *, fill=CARD, line=LINE, radius=True):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        PInches(left),
        PInches(top),
        PInches(width),
        PInches(height),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = PPt(1.0)
    if radius:
        try:
            shp.adjustments[0] = 0.06
        except Exception:
            pass
    return shp


def add_bg(slide):
    add_rect(slide, 0, 0, 13.333, 7.5, fill=PAPER, line=None, radius=False)


def add_footer(slide, page):
    add_text(
        slide,
        0.5,
        7.12,
        12.3,
        0.28,
        f"TransferBO · 共享条件库历史标签迁移 · {page:02d} / {TOTAL:02d}",
        size=10,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def header(slide, kicker: str, title: str):
    add_rect(slide, 0.55, 0.38, 0.1, 0.7, fill=ACCENT, line=None, radius=False)
    add_text(slide, 0.8, 0.28, 11.5, 0.28, kicker, size=12, color=MUTED, bold_first=True)
    add_text(slide, 0.8, 0.55, 11.8, 0.5, title, size=22, color=INK, bold_first=True)


def add_pic(slide, path: Path, left, top, width=None, height=None):
    if not path.exists():
        add_text(slide, left, top, 4, 0.4, f"[缺图] {path.name}", size=11, color=ACCENT2)
        return None
    kw = {}
    if width is not None:
        kw["width"] = PInches(width)
    if height is not None:
        kw["height"] = PInches(height)
    return slide.shapes.add_picture(str(path), PInches(left), PInches(top), **kw)


def build_ppt() -> Path:
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)
    blank = prs.slide_layouts[6]
    page = 0

    def new():
        nonlocal page
        page += 1
        s = prs.slides.add_slide(blank)
        add_bg(s)
        return s

    def section_chip(slide, label, left=0.7, top=1.35):
        add_rect(slide, left, top, 1.55, 0.32, fill=ACCENT, line=None, radius=False)
        add_text(slide, left + 0.08, top + 0.02, 1.4, 0.28, label, size=11, color=CARD, bold_first=True)

    # ---------- 01 封面 ----------
    s = new()
    add_rect(s, 0, 0, 13.333, 7.5, fill=PRGB(0xF0, 0xEB, 0xE3), line=None, radius=False)
    add_rect(s, 0, 0, 0.22, 7.5, fill=ACCENT, line=None, radius=False)
    add_text(s, 0.9, 1.35, 11.5, 0.35, "研究进展汇报 · 实验全链路", size=14, color=MUTED, bold_first=True)
    add_text(
        s,
        0.9,
        1.9,
        11.5,
        1.5,
        ["共享反应条件库上的历史标签池化", "相对冷启动不是安全默认策略"],
        size=28,
        color=INK,
        bold_first=True,
        space_after=8,
    )
    add_text(
        s,
        0.9,
        3.7,
        11.5,
        1.8,
        [
            "结构：实验设计 → 实验目的 → 实验结果 → 数据分析",
            "验证链：C1（多表示）→ S0（匹配 init）→ A2（秩池化）→ A3（源降权）+ 边界库对照",
            "阶段：主验证闭环 + 第二家族边界（amination 全量）· 2026-08",
        ],
        size=15,
        color=MUTED,
        space_after=6,
    )
    add_footer(s, page)

    # ---------- 02 目录 ----------
    s = new()
    header(s, "00  目录", "本次汇报结构")
    items = [
        ("一、问题与目的", "科学问题、总体实验目的、阶段结论边界"),
        ("二、数据与设计", "库来源/结构 · 协议 · 方法阶梯 · 分析协议"),
        ("三、实验结果", "C1 · S0 · A2 · A3 · amination 主数字"),
        ("四、数据分析", "预算窗口 · 异质性 · A3 sanity · 假设检验"),
        ("五、边界与收束", "Doyle / PK2022 / amination · 结论 · 待讨论"),
    ]
    y = 1.55
    for title, body in items:
        add_rect(s, 0.7, y, 11.9, 0.9, fill=CARD, line=LINE)
        add_text(s, 0.95, y + 0.18, 3.2, 0.55, title, size=16, color=ACCENT, bold_first=True, valign=MSO_ANCHOR.MIDDLE)
        add_text(s, 4.3, y + 0.18, 8.0, 0.55, body, size=15, color=INK, valign=MSO_ANCHOR.MIDDLE)
        y += 1.0
    add_footer(s, page)

    # ---------- 03 问题与总体目的 ----------
    s = new()
    header(s, "01  问题与目的", "我们锁定回答什么、为什么做这一套实验")
    section_chip(s, "目的")
    add_rect(s, 0.7, 1.8, 11.9, 1.55, fill=CARD, line=LINE)
    add_text(
        s,
        0.95,
        2.0,
        11.4,
        1.2,
        [
            "可操作问题：相对目标板冷启动，无 task ID、持续并入历史标签的池化，是不是安全默认？",
            "不是：历史数据有没有用 / 一切迁移是否无效 / MTGP 会不会成功。",
        ],
        size=16,
        color=INK,
        space_after=8,
    )
    cards = [
        (0.7, "总体目的", "在共享条件库、跨反应变体\n设定下，量化无任务身份\n标签池化的利 / 零 / 害"),
        (4.7, "验证布局", "多因子条件库为主验证\n多表示 + 稳健性阶梯\n外加边界库对照"),
        (8.7, "阶段结论", "不是安全默认\n解读：consistent with\ntask mismatch\n（非唯一机制断言）"),
    ]
    for left, title, body in cards:
        add_rect(s, left, 3.6, 3.7, 2.85, fill=CARD, line=LINE)
        add_rect(s, left, 3.6, 3.7, 0.45, fill=ACCENT, line=None, radius=False)
        add_text(s, left + 0.2, 3.65, 3.3, 0.35, title, size=14, color=CARD, bold_first=True)
        add_text(s, left + 0.2, 4.25, 3.3, 2.0, body.split("\n"), size=14, color=INK, space_after=4)
    add_footer(s, page)

    # ---------- 03b 数据来源与结构 ----------
    s = new()
    header(s, "02  数据", "本工作用到的共享条件库：来源与结构")
    section_chip(s, "设计")
    add_text(
        s,
        0.7,
        1.5,
        11.9,
        0.35,
        "共同点：离散候选集 X 跨任务共享；任务身份（底物/变体）不写入 X。差别在 X 的因子维数与化学问题。",
        size=12,
        color=MUTED,
    )
    # header row
    add_rect(s, 0.55, 1.95, 12.2, 0.42, fill=NAVY, line=None, radius=False)
    headers = [
        (0.65, 2.0, "数据集", 2.0),
        (2.7, 2.0, "来源", 2.6),
        (5.4, 2.0, "共享 X", 2.4),
        (7.9, 2.0, "任务 =", 2.0),
        (10.0, 2.0, "规模 / 角色", 2.5),
    ]
    for left, top, txt, w in headers:
        add_text(s, left, top, w, 0.3, txt, size=11, color=CARD, bold_first=True)
    rows = [
        (
            "EDBO Suzuki",
            "Shields et al.\nNature 2021",
            "L×B×Sol\n≈308 条件",
            "底物对\n(E + N)",
            "8 靶 · 主验证",
        ),
        (
            "EDBO amination",
            "同 EDBO2021\naryl_amination",
            "Add×Base×Lig\n≈260 条件",
            "芳基卤\n底物",
            "8 板 · 第二家族",
        ),
        (
            "Doyle2018",
            "Ahneman et al.\nScience 2018\n(BH CN HTE)",
            "L×B×Add\n≈240 条件",
            "芳基卤\n底物",
            "15×240 · 外部边界",
        ),
        (
            "PK2022",
            "Prieto–Kullmer\nScience 2022\n(CHAOS 四板)",
            "添加剂一维\n≈720 / 板",
            "反应变体\nplate_1–4",
            "4×720 · 1D 边界",
        ),
    ]
    y = 2.45
    for name, src, x, task, role in rows:
        add_rect(s, 0.55, y, 12.2, 1.05, fill=CARD, line=LINE)
        add_text(s, 0.65, y + 0.2, 2.0, 0.7, name, size=12, color=ACCENT, bold_first=True)
        add_text(s, 2.7, y + 0.12, 2.6, 0.85, src.split("\n"), size=11, color=INK, space_after=1)
        add_text(s, 5.4, y + 0.12, 2.4, 0.85, x.split("\n"), size=11, color=INK, space_after=1)
        add_text(s, 7.9, y + 0.12, 2.0, 0.85, task.split("\n"), size=11, color=INK, space_after=1)
        add_text(s, 10.0, y + 0.2, 2.5, 0.7, role, size=11, color=MUTED)
        y += 1.1
    add_footer(s, page)

    # ---------- 04 实验设计：体系与协议 ----------
    s = new()
    header(s, "03  实验设计", "体系设定与锁定协议")
    section_chip(s, "设计")
    add_pic(s, FIGS / "fig1_same_library_transfer_schematic.png", 0.55, 1.75, width=6.2)
    add_rect(s, 7.0, 1.75, 5.7, 4.85, fill=CARD, line=LINE)
    add_text(s, 7.25, 1.95, 5.2, 0.35, "协议常数（冻结）", size=13, color=MUTED, bold_first=True)
    add_text(
        s,
        7.25,
        2.4,
        5.2,
        3.9,
        [
            "• 设定：共享离散条件库 X",
            "• X = 配体 × 碱 × 溶剂（多因子）",
            "• 任务 = 不同底物 / 反应变体",
            "  （底物身份不写入 X）",
            "• 推断：有向 source→target pair",
            "• 代理：GP Matérn(ν=2.5)+White",
            "• 采集：EI；normalize_y=True",
            "• n_init=20；B=100 目标查询",
            "• 源标签不计入目标预算",
            "• seeds = 0…19；无 task ID",
            "• 主验证库规模：308 条件 × 8 板",
        ],
        size=13,
        color=INK,
        space_after=3,
    )
    add_footer(s, page)

    # ---------- 05 方法阶梯与各臂目的 ----------
    s = new()
    header(s, "03  实验设计", "方法阶梯 A0–A3：每臂在检验什么")
    section_chip(s, "设计+目的")
    rows = [
        ("A0 冷启动", "不用历史", "建立强目标板基线", "对照臂"),
        ("A1 原始池化", "源 y 直接混入同一 GP", "产率是否可交换？", "主效应 C1"),
        ("S0 匹配 init", "冷/池化共享目标 init", "负迁移是否只是 init 错配？", "稳健性"),
        ("A2 秩池化", "任务内百分位后再池化", "是否主要是标度问题？", "机制排除"),
        ("A3 源降权", "α_src=1e-4/w_s", "中等可靠性收缩是否够？", "简单补救"),
    ]
    add_rect(s, 0.7, 1.75, 11.9, 0.45, fill=ACCENT, line=None, radius=False)
    add_text(s, 0.9, 1.82, 2.4, 0.35, "臂", size=12, color=CARD, bold_first=True)
    add_text(s, 3.3, 1.82, 3.0, 0.35, "做法", size=12, color=CARD, bold_first=True)
    add_text(s, 6.5, 1.82, 3.6, 0.35, "实验目的（检验假设）", size=12, color=CARD, bold_first=True)
    add_text(s, 10.2, 1.82, 2.1, 0.35, "角色", size=12, color=CARD, bold_first=True)
    y = 2.3
    for name, how, purpose, role in rows:
        add_rect(s, 0.7, y, 11.9, 0.82, fill=CARD, line=LINE)
        add_text(s, 0.9, y + 0.2, 2.4, 0.5, name, size=14, color=ACCENT, bold_first=True)
        add_text(s, 3.3, y + 0.2, 3.0, 0.5, how, size=13, color=INK)
        add_text(s, 6.5, y + 0.2, 3.6, 0.5, purpose, size=13, color=INK)
        add_text(s, 10.2, y + 0.2, 2.1, 0.5, role, size=13, color=MUTED)
        y += 0.88
    add_footer(s, page)

    # ---------- 06 计算网格 ----------
    s = new()
    header(s, "04  实验设计", "计算网格与结果库存（验证全貌）")
    section_chip(s, "设计")
    grids = [
        ("C1 多表示", "7200 JSON", "cold+label × Morgan/DRFP/DFT", "主效应：是否表示特异"),
        ("S0 匹配 init", "2560 JSON", "多表示；init 100% 匹配", "排除初始化伪影"),
        ("A2 秩池化", "2240 JSON", "百分位秩后再池化", "排除纯标度解释"),
        ("A3 源降权", "6720 JSON", "w∈{0.1,0.25,0.5}×多表示", "检验简单可靠性收缩"),
    ]
    y = 1.75
    for name, n, detail, why in grids:
        add_rect(s, 0.7, y, 11.9, 1.05, fill=CARD, line=LINE)
        add_text(s, 0.95, y + 0.15, 2.3, 0.7, name, size=16, color=ACCENT, bold_first=True, valign=MSO_ANCHOR.MIDDLE)
        add_text(s, 3.4, y + 0.15, 2.0, 0.7, n, size=15, color=ACCENT2, bold_first=True, valign=MSO_ANCHOR.MIDDLE)
        add_text(s, 5.5, y + 0.15, 4.0, 0.7, detail, size=13, color=INK, valign=MSO_ANCHOR.MIDDLE)
        add_text(s, 9.6, y + 0.15, 2.7, 0.7, why, size=13, color=MUTED, valign=MSO_ANCHOR.MIDDLE)
        y += 1.15
    add_text(
        s,
        0.7,
        6.5,
        11.9,
        0.4,
        "边界 / 补充验证：Doyle2018（OHE≈2400）· PK2022/CHAOS（≈840）· xTB pilot · EDBO amination 全量 2560",
        size=12,
        color=MUTED,
    )
    add_footer(s, page)

    # ---------- 07 数据分析协议 ----------
    s = new()
    header(s, "05  实验设计", "数据分析协议")
    section_chip(s, "分析协议")
    boxes = [
        (0.7, "指标", "frac(B)=已见最佳产率/板上最优\nΔfrac = label − cold\n主报 pair 平均 Δfrac"),
        (4.7, "推断单位", "有向 source→target pair\nN≈56；seed 估算法波动\n禁把 1120 轨迹当 IID"),
        (8.7, "预算窗口", "主窗口 B=30 / 40 / 50\nB=100 = 天花板对照\n中段 headroom 仍在"),
    ]
    for left, title, body in boxes:
        add_rect(s, left, 1.8, 3.7, 2.7, fill=CARD, line=LINE)
        add_rect(s, left, 1.8, 3.7, 0.45, fill=NAVY, line=None, radius=False)
        add_text(s, left + 0.2, 1.85, 3.3, 0.35, title, size=14, color=CARD, bold_first=True)
        add_text(s, left + 0.2, 2.45, 3.3, 1.9, body.split("\n"), size=14, color=INK, space_after=5)
    add_rect(s, 0.7, 4.75, 11.9, 1.75, fill=CARD, line=LINE)
    add_text(s, 0.95, 4.9, 11.4, 0.35, "分析规则（冻结）", size=13, color=MUTED, bold_first=True)
    add_text(
        s,
        0.95,
        5.35,
        11.4,
        1.0,
        [
            "• CI：对 pair 做 bootstrap（非对轨迹）；同时报 n_pos / n_neg / n_near0",
            "• 主网格 init 0% 匹配 → 必须以 S0 复核；A2/A3 相对 S0 cold 比较",
            "• 不因追求正向结果改协议或挑选 pair；数字以汇总表为准",
        ],
        size=14,
        color=INK,
        space_after=4,
    )
    add_footer(s, page)

    # ---------- 08 结果 C1 ----------
    s = new()
    header(s, "06  实验结果 · C1", "原始池化 vs 冷启动（三表示）")
    section_chip(s, "结果")
    add_text(
        s,
        2.4,
        1.35,
        10.0,
        0.3,
        "目的：在锁定协议下，量化 task-agnostic 原始标签池化的平均效应",
        size=12,
        color=MUTED,
    )
    add_pic(s, FIGS / "fig_edbo_suzuki_C1_pair_delta_by_budget.png", 0.45, 1.75, width=7.5)
    add_rect(s, 8.15, 1.75, 4.55, 4.75, fill=CARD, line=LINE)
    add_text(s, 8.4, 1.95, 4.1, 0.35, "pair Δfrac（headline）", size=13, color=MUTED, bold_first=True)
    add_text(
        s,
        8.4,
        2.4,
        4.1,
        3.8,
        [
            "B=40",
            "  Morgan  −0.041",
            "  DRFP    −0.036",
            "  DFT     −0.036",
            "",
            "B=30 同向略负",
            "B=100 压缩至 ≈−0.01",
            "",
            "95% pair-bootstrap CI",
            "中段多不含 0",
            "表示间相关高（r≳0.86）",
        ],
        size=14,
        color=INK,
        space_after=3,
    )
    add_footer(s, page)

    # ---------- 09 数据分析 C1 ----------
    s = new()
    header(s, "07  数据分析 · C1", "为何主窗口是 B=30–50，而非 B=100")
    section_chip(s, "分析")
    add_pic(s, ESI / "fig_edbo_suzuki_headroom_vs_delta_frac_B40.png", 0.5, 1.75, width=6.6)
    add_rect(s, 7.35, 1.75, 5.35, 4.75, fill=CARD, line=LINE)
    add_text(s, 7.6, 1.95, 4.9, 0.35, "读图要点", size=13, color=MUTED, bold_first=True)
    add_text(
        s,
        7.6,
        2.45,
        4.9,
        3.8,
        [
            "1. 三表示同向 → 不是指纹特例",
            "2. 中段伤害最大 → 仍有 headroom",
            "3. B=100 冷启动接近天花板",
            "   → 终值 Δfrac 被压缩",
            "4. 故「接近零@100」≠",
            "   「池化无害」",
            "",
            "分析结论：",
            "原始池化在中段平均略害；",
            "不是安全默认。",
        ],
        size=14,
        color=INK,
        space_after=4,
    )
    add_footer(s, page)

    # ---------- 10 结果 S0 ----------
    s = new()
    header(s, "08  实验结果 · S0", "匹配目标板初始化（稳健性）")
    section_chip(s, "结果")
    add_text(
        s,
        2.4,
        1.35,
        10.0,
        0.3,
        "目的：排除「主网格 cold↔label init 0% 匹配」造成的伪负迁移",
        size=12,
        color=MUTED,
    )
    add_pic(s, FIGS / "fig_edbo_suzuki_s0_vs_main_pair_delta.png", 0.45, 1.75, width=7.6)
    add_rect(s, 8.25, 1.75, 4.45, 4.75, fill=CARD, line=LINE)
    add_text(
        s,
        8.5,
        2.05,
        4.0,
        4.2,
        [
            "设计：init 先采目标板；",
            "源采样 RNG = seed+1000003",
            "→ 同 (target,seed) 匹配 100%",
            "",
            "B=40 pair Δfrac",
            "Morgan  Main −0.041",
            "        S0   −0.034",
            "DFT     Main −0.036",
            "        S0   −0.030",
            "",
            "分析：略减负（≈0.006），",
            "符号不翻；CI 仍多不含 0。",
            "init 错配会放大，非主因。",
        ],
        size=13,
        color=INK,
        space_after=3,
    )
    add_footer(s, page)

    # ---------- 11 结果 A2 ----------
    s = new()
    header(s, "09  实验结果 · A2", "任务内秩 / 百分位池化")
    section_chip(s, "结果")
    add_rect(s, 0.7, 1.8, 11.9, 1.2, fill=CARD, line=LINE)
    add_text(
        s,
        0.95,
        2.0,
        11.4,
        0.85,
        [
            "目的：若伤害主要来自不同底物产率标度不可比，则秩变换应明显修复。",
            "设计：各任务内将 y 转为百分位秩，再在秩空间做与 A1 相同的无 task ID 池化；对照 = S0 cold。",
        ],
        size=15,
        color=INK,
        space_after=5,
    )
    add_rect(s, 0.7, 3.25, 5.85, 3.25, fill=CARD, line=LINE)
    add_rect(s, 6.8, 3.25, 5.85, 3.25, fill=CARD, line=LINE)
    add_text(s, 0.95, 3.45, 5.4, 0.35, "B=40 结果（vs S0 cold）", size=13, color=MUTED, bold_first=True)
    add_text(
        s,
        0.95,
        3.95,
        5.4,
        2.3,
        [
            "A1 raw   Morgan −0.034 / DFT −0.030",
            "A2 rank  Morgan −0.026 / DFT −0.029",
            "",
            "略减负，未翻正；DFT 几乎不动。",
        ],
        size=15,
        color=INK,
        space_after=5,
    )
    add_text(s, 7.05, 3.45, 5.4, 0.35, "分析解读", size=13, color=MUTED, bold_first=True)
    add_text(
        s,
        7.05,
        3.95,
        5.4,
        2.3,
        [
            "排除「只是产量标度」的单一解释。",
            "更符合景观失配 / 任务不可交换：",
            "即便秩对齐，共享核仍把",
            "源最优当目标信号。",
        ],
        size=15,
        color=INK,
        space_after=5,
    )
    add_footer(s, page)

    # ---------- 12 结果 A3 + ladder ----------
    s = new()
    header(s, "10  实验结果 · A3", "源标签噪声降权（方法阶梯总图）")
    section_chip(s, "结果")
    add_text(
        s,
        2.4,
        1.35,
        10.0,
        0.3,
        "目的：中等源降权能否在无 task ID 设定下稳定扭转负迁移？",
        size=12,
        color=MUTED,
    )
    add_pic(s, FIGS / "fig_edbo_suzuki_ladder_A1A2A3_B40.png", 0.4, 1.7, width=8.0)
    add_rect(s, 8.55, 1.7, 4.2, 4.8, fill=CARD, line=LINE)
    add_text(
        s,
        8.8,
        1.9,
        3.8,
        4.4,
        [
            "设计：α_src=1e-4/w_s",
            "w∈{0.1, 0.25, 0.5}",
            "对照：S0 cold / A1",
            "",
            "B=40 vs S0 cold",
            "A1      ≈ −0.034/−0.030",
            "A3 w0.1 ≈ −0.036/−0.030",
            "A3 w0.25≈ −0.036/−0.031",
            "A3 w0.5 ≈ −0.036/−0.031",
            "",
            "与 A1 几乎重合。",
            "不是「降权永远无效」，",
            "而是测试档未修复。",
        ],
        size=13,
        color=INK,
        space_after=3,
    )
    add_footer(s, page)

    # ---------- 13 A3 sanity 分析 ----------
    s = new()
    header(s, "11  数据分析 · A3", "健全性：权重进了模型，但采集路径常不变")
    section_chip(s, "分析")
    cols = [
        (0.7, "实现核对", "LabelWeight → α_src\n→ sklearn GPR alpha\n接线真实，非空转"),
        (4.7, "后验探针", "固定 init：\nmax|Δμ|(w=0.1 vs 1)\n≈ 0.05（输出尺度）\nw=1e-4 才明显近 cold"),
        (8.7, "轨迹统计", "Morgan：~78% 路径\n在 w=0.1–0.5 间相同\nDFT：~73%\n极端小 w → 趋 cold"),
    ]
    for left, title, body in cols:
        add_rect(s, left, 1.8, 3.7, 3.0, fill=CARD, line=LINE)
        add_rect(s, left, 1.8, 3.7, 0.45, fill=ACCENT2, line=None, radius=False)
        add_text(s, left + 0.2, 1.85, 3.3, 0.35, title, size=14, color=CARD, bold_first=True)
        add_text(s, left + 0.2, 2.45, 3.3, 2.2, body.split("\n"), size=14, color=INK, space_after=4)
    add_rect(s, 0.7, 5.05, 11.9, 1.45, fill=CARD, line=LINE)
    add_text(
        s,
        0.95,
        5.25,
        11.4,
        1.1,
        [
            "结论口径：中等噪声膨胀未能一致优于原始池化 / 扭转负迁移；早期 n_s≫n_t 时源点仍被高度信任。",
            "避免过度解读：不是「源降权永远无效」；也不把 w=1e-4（≈退回冷启动）当作有效方法臂。",
        ],
        size=14,
        color=INK,
        space_after=5,
    )
    add_footer(s, page)

    # ---------- 14 异质性 ----------
    s = new()
    header(s, "12  数据分析 · 异质性", "pair 层面：平均略负 + 厚负尾")
    section_chip(s, "分析")
    add_pic(s, ESI / "fig_edbo_suzuki_morgan_pair_delta_frac_heatmap_B40.png", 0.4, 1.7, width=6.5)
    add_rect(s, 7.15, 1.7, 5.55, 4.8, fill=CARD, line=LINE)
    add_text(s, 7.4, 1.9, 5.1, 0.35, "分析要点（Morgan @ B=40）", size=13, color=MUTED, bold_first=True)
    add_text(
        s,
        7.4,
        2.4,
        5.1,
        3.9,
        [
            "• 效应高度 pair 依赖",
            "• 正尾薄、负尾厚",
            "• 八个目标的目标聚合",
            "  中段均值均 ≤ 0",
            "• 低天花板目标（如 t12）",
            "  平均伤害更大",
            "",
            "含义：不是「一个坏 pair」",
            "也不是「平均正、偶发负」。",
            "实践上不能指望「换个源",
            "就默认变好」。",
        ],
        size=14,
        color=INK,
        space_after=3,
    )
    add_footer(s, page)

    # ---------- 15 假设检验综合表 ----------
    s = new()
    header(s, "13  数据分析 · 综合", "对照假设 → 结果 → 排除什么")
    section_chip(s, "分析")
    add_rect(s, 0.7, 1.75, 11.9, 0.42, fill=NAVY, line=None, radius=False)
    add_text(s, 0.9, 1.82, 2.6, 0.3, "对照", size=12, color=CARD, bold_first=True)
    add_text(s, 3.6, 1.82, 3.4, 0.3, "检验假设", size=12, color=CARD, bold_first=True)
    add_text(s, 7.1, 1.82, 2.4, 0.3, "结果", size=12, color=CARD, bold_first=True)
    add_text(s, 9.6, 1.82, 2.8, 0.3, "排除", size=12, color=CARD, bold_first=True)
    rows = [
        ("A0 vs A1", "源/目标 y 可交换", "中段平均负迁移", "原始标签可交换"),
        ("S0", "伤害仅来自 init 错配", "仍负，略减轻", "init 为唯一原因"),
        ("A1 vs A2", "主要是产率标度", "略减负，未翻正", "纯标度故事"),
        ("A1 vs A3", "中等可靠性收缩够用", "≈A1，未稳定修复", "简单降权即够"),
    ]
    y = 2.3
    for a, b, c, d in rows:
        add_rect(s, 0.7, y, 11.9, 0.95, fill=CARD, line=LINE)
        add_text(s, 0.9, y + 0.25, 2.6, 0.5, a, size=14, color=ACCENT, bold_first=True)
        add_text(s, 3.6, y + 0.25, 3.4, 0.5, b, size=13, color=INK)
        add_text(s, 7.1, y + 0.25, 2.4, 0.5, c, size=13, color=INK)
        add_text(s, 9.6, y + 0.25, 2.8, 0.5, d, size=13, color=MUTED)
        y += 1.05
    add_text(
        s,
        0.7,
        6.55,
        11.9,
        0.35,
        "综合解读（非唯一证明）：与 task mismatch 一致 —— 共享 X ⇏ 跨底物历史 y 可交换。",
        size=13,
        color=ACCENT,
        bold_first=True,
    )
    add_footer(s, page)

    # ---------- 16 边界 ----------
    s = new()
    header(s, "14  边界对照", "设计角色：说明「并非普适禁止迁移」")
    section_chip(s, "设计+结果")
    boxes = [
        (
            0.7,
            "Doyle2018",
            "Ahneman Science 2018\nBH CN HTE；X=L×B×Add≈240\n任务=芳基卤底物（15×240）\nlabel(OHE) pair Δ≈+0.064\n→ 外部边界：别处可正",
        ),
        (
            4.7,
            "PK2022",
            "Prieto–Kullmer Science 2022\nCHAOS 四板；X≈添加剂720\n任务=反应变体 plate_1–4\nlabel pair Δ≈+0.13~+0.15\n→ 1D 边界：不可当主证据",
        ),
        (
            8.7,
            "EDBO amination",
            "同 EDBO2021 第二家族\n全量 2560 已完成\nDFT +0.043；Morgan ≈0\nfamily-dependent\n→ 见下一页",
        ),
    ]
    for left, title, body in boxes:
        add_rect(s, left, 1.8, 3.7, 4.6, fill=CARD, line=LINE)
        add_rect(s, left, 1.8, 3.7, 0.5, fill=NAVY, line=None, radius=False)
        add_text(s, left + 0.2, 1.9, 3.3, 0.35, title, size=15, color=CARD, bold_first=True)
        add_text(s, left + 0.25, 2.55, 3.2, 3.5, body.split("\n"), size=14, color=INK, space_after=7)
    add_footer(s, page)

    # ---------- 16b amination W8 ----------
    s = new()
    header(s, "14b  第二家族", "EDBO amination min S0：family-dependent")
    section_chip(s, "结果+分析")
    add_text(
        s,
        0.7,
        1.55,
        11.9,
        0.35,
        "协议对齐 Suzuki S0：cold vs label_warm · Morgan+DFT · 8 底物 · 20 seeds · 2560 JSON · 主终点 mean Δfrac @ B∈{30,40,50}",
        size=12,
        color=MUTED,
    )
    add_rect(s, 0.7, 2.0, 5.9, 2.55, fill=CARD, line=LINE)
    add_rect(s, 6.85, 2.0, 5.75, 2.55, fill=CARD, line=LINE)
    add_text(s, 0.95, 2.15, 5.4, 0.35, "Amination（pair 平均）", size=13, color=MUTED, bold_first=True)
    add_text(
        s,
        0.95,
        2.6,
        5.4,
        1.8,
        [
            "• DFT：mean +0.043（59% >+0.02）",
            "• Morgan：mean +0.013（50% near0）",
            "• 相对 Suzuki C1 ≈ −0.03x：符号可相反",
        ],
        size=14,
        color=INK,
        space_after=5,
    )
    add_text(s, 7.1, 2.15, 5.3, 0.35, "异质性（关键）", size=13, color=MUTED, bold_first=True)
    add_text(
        s,
        7.1,
        2.6,
        5.3,
        1.8,
        [
            "• 正均值被 target=sub_s4 拉高",
            "• 去掉 s4：DFT ≈ +0.020；Morgan ≈ 0",
            "• 不是八个底物普遍受益",
        ],
        size=14,
        color=INK,
        space_after=5,
    )
    add_rect(s, 0.7, 4.8, 11.9, 1.7, fill=CARD, line=LINE)
    add_text(s, 0.95, 5.0, 11.4, 0.35, "写作口径（已锁定）", size=13, color=ACCENT, bold_first=True)
    add_text(
        s,
        0.95,
        5.45,
        11.4,
        0.9,
        [
            "不做跨家族「一致不安全默认」升级。Amination = 边界：同协议下效应可 family-dependent。",
            "主主张仍锚定多因子 Suzuki：naive 持续池化不是安全默认；amination 说明不能外推成「处处负」。",
        ],
        size=14,
        color=INK,
        space_after=4,
    )
    add_footer(s, page)

    # ---------- 17 总结论 ----------
    s = new()
    header(s, "15  收束", "阶段结论与解读边界")
    add_rect(s, 0.7, 1.55, 11.9, 1.7, fill=CARD, line=LINE)
    add_text(
        s,
        0.95,
        1.75,
        11.4,
        1.35,
        [
            "在主验证多因子共享库上：无 task ID 的持续历史标签池化，相对冷启动不是安全默认。",
            "多表示同向负；匹配 init / 秩 / 名义降权未翻正。第二家族 amination 平均偏正/近零 → 效应 family-dependent。",
        ],
        size=15,
        color=INK,
        space_after=6,
    )
    add_rect(s, 0.7, 3.5, 5.9, 3.0, fill=CARD, line=LINE)
    add_rect(s, 6.85, 3.5, 5.75, 3.0, fill=CARD, line=LINE)
    add_text(s, 0.95, 3.7, 5.4, 0.35, "当前支持的解读", size=13, color=ACCENT, bold_first=True)
    add_text(
        s,
        0.95,
        4.2,
        5.4,
        2.1,
        [
            "• 共享 X ≠ 历史 y 可默认混用",
            "• 与 task mismatch 相一致",
            "• 别处/他家族可正 → 非普适禁令",
        ],
        size=14,
        color=INK,
        space_after=5,
    )
    add_text(s, 7.1, 3.7, 5.3, 0.35, "不宜外推", size=13, color=ACCENT2, bold_first=True)
    add_text(
        s,
        7.1,
        4.2,
        5.3,
        2.1,
        [
            "• 历史数据没用 / 一切 transfer 无效",
            "• task-aware / MTGP 也必然失败",
            "• 任意反应家族都必然负迁移",
        ],
        size=14,
        color=INK,
        space_after=5,
    )
    add_footer(s, page)

    # ---------- 18 进度 ----------
    s = new()
    header(s, "16  进度", "工作进展")
    add_rect(s, 0.7, 1.6, 5.9, 4.9, fill=CARD, line=LINE)
    add_text(s, 0.95, 1.8, 5.4, 0.4, "已完成", size=14, color=ACCENT, bold_first=True)
    add_text(
        s,
        0.95,
        2.35,
        5.4,
        3.9,
        [
            "• C1 / S0 / A2 / A3 全套验证",
            "• Task-level 重推断 + A2/A3 审计",
            "• 相似度机制（top-k ↔ Δfrac）",
            "• Amination 全量 2560（W8）",
            "• S5 / task-ID 策略臂 + 小试点",
        ],
        size=15,
        color=INK,
        space_after=6,
    )
    add_rect(s, 6.9, 1.6, 5.7, 4.9, fill=CARD, line=LINE)
    add_text(s, 7.15, 1.8, 5.2, 0.4, "可选下一步", size=14, color=ACCENT2, bold_first=True)
    add_text(
        s,
        7.15,
        2.35,
        5.2,
        3.9,
        [
            "• 扩跑 W7 nsweep / S5 / task-ID",
            "  （现为小试点）",
            "• 或推进更强 task-aware（MTGP）",
            "• xTB / 边界库多表示：低优先",
            "",
            "当前阶段结论：",
            "主库上 naive 池化非安全默认；",
            "跨家族不可一刀切。",
        ],
        size=14,
        color=INK,
        space_after=5,
    )
    add_footer(s, page)

    # ---------- 19 待讨论 ----------
    s = new()
    header(s, "17  待讨论", "需要确认的几点")
    qs = [
        "1. Amination 作为 family-dependent 边界写入正文/SI，是否同意不做跨家族强升级？",
        "2. 下一阶段是否扩跑 S5（仅 init）全网格，作为「然后怎么办」主对照？",
        "3. 轻量 task-ID 试点已≈A1；是否值得上更强 MTGP，还是先写清局限？",
        "4. 边界库是否需要多表示重跑，还是维持现状作对照即可？",
    ]
    y = 1.6
    for q in qs:
        add_rect(s, 0.7, y, 11.9, 1.1, fill=CARD, line=LINE)
        add_text(s, 0.95, y + 0.25, 11.4, 0.7, q, size=15, color=INK, valign=MSO_ANCHOR.MIDDLE)
        y += 1.25
    add_footer(s, page)

    # ---------- 20 结束 ----------
    s = new()
    add_rect(s, 0, 0, 13.333, 7.5, fill=PRGB(0xF0, 0xEB, 0xE3), line=None, radius=False)
    add_rect(s, 0, 0, 0.22, 7.5, fill=ACCENT, line=None, radius=False)
    add_text(s, 0.9, 2.2, 11.5, 0.5, "谢谢，欢迎讨论与指正", size=28, color=INK, bold_first=True)
    add_text(
        s,
        0.9,
        3.2,
        11.5,
        2.0,
        [
            "汇报材料：docs/briefings/（v0.7）",
            "数字汇总：results/paper_stats/EXPERIMENT_SUMMARY.md",
            "Amination：edbo_amination_min_s0_SUMMARY.md",
            "结果目录：results/external_* · results/transfer_grid*",
        ],
        size=15,
        color=MUTED,
        space_after=6,
    )
    add_footer(s, page)

    if page != TOTAL:
        # keep footer consistent if slide count drifts
        pass

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(PPT_OUT)
    try:
        import shutil

        shutil.copy2(PPT_OUT, PPT_ALIAS)
    except Exception:
        pass
    return PPT_OUT


# ----------------- DOCX -----------------


def set_run_font(run, *, east_asia="宋体", ascii_font="Times New Roman", size_pt=12):
    run.font.name = ascii_font
    run.font.size = Pt(size_pt)
    rPr = run._element.get_or_add_rPr()
    rFonts = rPr.get_or_add_rFonts()
    rFonts.set(qn("w:eastAsia"), east_asia)


def add_h(doc, text, level=1):
    p = doc.add_heading(text, level=level)
    for run in p.runs:
        set_run_font(run, east_asia="黑体", ascii_font="Arial", size_pt=16 if level == 1 else 13)


def add_p(doc, text, *, indent=True, bold=False, size=12, center=False):
    p = doc.add_paragraph()
    if center:
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    if indent:
        p.paragraph_format.first_line_indent = Cm(0.74)
    p.paragraph_format.line_spacing_rule = WD_LINE_SPACING.ONE_POINT_FIVE
    p.paragraph_format.space_after = Pt(6)
    run = p.add_run(text)
    run.bold = bold
    set_run_font(run, size_pt=size)
    return p


def add_bul(doc, text):
    p = doc.add_paragraph(style="List Bullet")
    p.paragraph_format.line_spacing_rule = WD_LINE_SPACING.ONE_POINT_FIVE
    run = p.add_run(text)
    set_run_font(run, size_pt=12)


def add_tbl(doc, headers, rows):
    t = doc.add_table(rows=1 + len(rows), cols=len(headers))
    t.style = "Table Grid"
    for j, h in enumerate(headers):
        cell = t.rows[0].cells[j]
        cell.text = ""
        run = cell.paragraphs[0].add_run(h)
        run.bold = True
        set_run_font(run, east_asia="黑体", size_pt=10.5)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = t.rows[i + 1].cells[j]
            cell.text = ""
            run = cell.paragraphs[0].add_run(val)
            set_run_font(run, size_pt=10.5)
    doc.add_paragraph()


def add_fig(doc, path: Path, caption: str):
    if not path.exists():
        add_p(doc, f"[缺图：{path.name}]", indent=False, bold=True)
        return
    doc.add_picture(str(path), width=Inches(5.9))
    doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
    cap = doc.add_paragraph()
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = cap.add_run(caption)
    set_run_font(run, east_asia="楷体", size_pt=10.5)


def build_docx() -> Path:
    """DOCX mirrors the advisor PPT: design → purpose → results → analysis."""
    doc = Document()
    for sec in doc.sections:
        sec.top_margin = Cm(2.4)
        sec.bottom_margin = Cm(2.4)
        sec.left_margin = Cm(2.6)
        sec.right_margin = Cm(2.6)

    t = doc.add_paragraph()
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = t.add_run("工作进展汇报")
    r.bold = True
    set_run_font(r, east_asia="黑体", ascii_font="Arial", size_pt=18)

    st = doc.add_paragraph()
    st.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = st.add_run(
        "共享反应条件库上的历史标签池化\n"
        "相对冷启动不是安全默认策略"
    )
    r.bold = True
    set_run_font(r, east_asia="黑体", ascii_font="Arial", size_pt=14)

    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = meta.add_run(
        "项目：TransferBO · 2026-08 · v0.7（含 amination 全量）\n"
        "结构与 PPT 对齐：实验设计 → 实验目的 → 实验结果 → 数据分析"
    )
    set_run_font(r, size_pt=10.5)
    r.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    add_h(doc, "一、问题与目的", 1)
    add_p(
        doc,
        "同一套离散反应条件库更换底物或反应变体时，历史产率要不要直接喂给新任务的贝叶斯优化？"
        "共享候选集 X 并不意味着共享响应函数 y(x)。",
    )
    add_p(
        doc,
        "可操作问题：相对目标板冷启动，无 task ID、持续并入历史标签的池化，是不是安全默认？",
        bold=True,
    )
    add_p(doc, "不是在问：历史数据有没有用、一切迁移是否无效、或多任务 GP 会不会成功。", indent=True)
    add_bul(doc, "总体目的：在共享条件库、跨反应变体设定下，量化无任务身份标签池化的利 / 零 / 害。")
    add_bul(doc, "验证布局：多因子条件库为主验证；多表示 + 稳健性阶梯；外加边界库对照。")
    add_bul(doc, "阶段结论：不是安全默认；解读为 consistent with task mismatch（非唯一机制断言）。")

    add_h(doc, "二、数据来源与结构", 1)
    add_p(
        doc,
        "本工作使用四套公开高通量实验（HTE）共享条件库。共同点：离散候选集 X 在任务间共享，"
        "而任务身份（底物或反应变体）不写入 X。差别在于 X 的因子维数与所回答的化学问题。",
        indent=False,
    )
    add_tbl(
        doc,
        ["数据集", "来源", "共享 X", "任务定义", "规模", "角色"],
        [
            [
                "EDBO Suzuki",
                "Shields et al., Nature 2021",
                "配体×碱×溶剂 ≈308",
                "底物对（亲电体+亲核体）",
                "主网格用 8 靶（全表 12 板）",
                "主验证",
            ],
            [
                "EDBO amination",
                "同 EDBO2021 aryl_amination",
                "添加剂×碱×配体 ≈260",
                "芳基卤底物",
                "8 板 × 260",
                "第二反应家族边界",
            ],
            [
                "Doyle2018",
                "Ahneman et al., Science 2018（Buchwald–Hartwig CN HTE）",
                "配体×碱×添加剂 ≈240",
                "芳基卤底物",
                "15 × 240",
                "外部边界（OHE 等）",
            ],
            [
                "PK2022",
                "Prieto–Kullmer et al., Science 2022；CHAOS 复用四板",
                "添加剂一维 ≈720/板",
                "反应变体 plate_1–4",
                "4 × 720",
                "1D 边界 / SI",
            ],
        ],
    )
    add_p(
        doc,
        "层级锁定：EDBO Suzuki 为唯一主证据设定；Doyle2018 / PK2022 / EDBO amination 只作边界或对照，"
        "不与主网格做效应量硬对齐，也不把 PK2022 的一维添加剂库当作多因子条件库主证据。",
        indent=False,
    )
    add_bul(doc, "Doyle2018（label_warm, OHE）：pair 平均 Δfrac ≈ +0.064（CI 不含 0）→ 别处可正。")
    add_bul(doc, "PK2022/CHAOS（label_warm）：Morgan/DRFP pair 平均 Δfrac ≈ +0.13 / +0.15 → 标签常有用；diversity-only 不稳。")
    add_bul(doc, "备注：CHAOS 指 Ranković 等对 PK2022 板的复用/BO 工作流，不是库名本身。")

    add_h(doc, "三、实验设计", 1)
    add_h(doc, "3.1 体系与协议", 2)
    add_fig(
        doc,
        FIGS / "fig1_same_library_transfer_schematic.png",
        "图1　共享条件库上历史标签池化示意",
    )
    add_bul(doc, "设定：共享离散条件库 X；主验证为多因子（配体 × 碱 × 溶剂）。")
    add_bul(doc, "任务 = 不同底物 / 反应变体（底物身份不写入 X）；约 308 条件 × 8 板 → 56 有向 pair。")
    add_bul(doc, "代理：GP Matérn(ν=2.5) + White；采集 EI；normalize_y=True。")
    add_bul(doc, "n_init=20；B=100 目标查询；源标签不计入目标预算；seeds=0…19；无 task ID。")

    add_h(doc, "3.2 方法阶梯：每臂检验什么", 2)
    add_tbl(
        doc,
        ["臂", "做法", "实验目的（检验假设）", "角色"],
        [
            ["A0 冷启动", "不用历史", "建立强目标板基线", "对照"],
            ["A1 原始池化", "源 y 直接混入同一 GP", "产率是否可交换？", "主效应 C1"],
            ["S0 匹配 init", "冷/池化共享目标 init", "负迁移是否只是 init 错配？", "稳健性"],
            ["A2 秩池化", "任务内百分位后再池化", "是否主要是标度问题？", "机制排除"],
            ["A3 源降权", "α_src=1e-4/w_s", "中等可靠性收缩是否够？", "简单补救"],
        ],
    )

    add_h(doc, "3.3 计算网格", 2)
    add_tbl(
        doc,
        ["网格", "规模", "内容", "目的"],
        [
            ["C1 多表示", "7200", "cold+label × Morgan/DRFP/DFT", "主效应；是否表示特异"],
            ["S0 匹配 init", "2560", "多表示；init 100% 匹配", "排除初始化伪影"],
            ["A2 秩池化", "2240", "百分位秩后再池化", "排除纯标度解释"],
            ["A3 源降权", "6720", "w∈{0.1,0.25,0.5}×多表示", "检验简单可靠性收缩"],
        ],
    )
    add_p(
        doc,
        "边界 / 补充验证：Doyle2018（约 2400）、PK2022/CHAOS（约 840）、xTB pilot；"
        "第二反应家族 EDBO amination 全量 2560（已完成）。",
        indent=False,
    )

    add_h(doc, "3.4 数据分析协议", 2)
    add_bul(doc, "指标：frac(B)=已见最佳产率/板上最优；Δfrac = label − cold；主报 pair 平均 Δfrac。")
    add_bul(doc, "推断单位：有向 source→target pair（N≈56）；seed 估算法波动；禁止把轨迹当 IID。")
    add_bul(doc, "预算窗口：主窗口 B=30/40/50；B=100 为天花板对照（中段仍有 headroom）。")
    add_bul(doc, "CI：对 pair 做 bootstrap；同时报 n_pos / n_neg / n_near0。")
    add_bul(doc, "主网格曾存在 init 不匹配 → 必须以 S0 复核；A2/A3 相对 S0 cold 比较。")
    add_bul(doc, "不因追求正向结果改协议或挑选 pair。")

    add_h(doc, "四、实验结果", 1)
    add_h(doc, "4.1 C1：原始池化 vs 冷启动（多表示）", 2)
    add_p(doc, "目的：在锁定协议下，量化无任务身份原始标签池化的平均效应。", indent=False)
    add_fig(
        doc,
        FIGS / "fig_edbo_suzuki_C1_pair_delta_by_budget.png",
        "图2　C1：多表示 pair Δfrac 随预算",
    )
    add_tbl(
        doc,
        ["表示", "B=40 Δfrac", "备注"],
        [
            ["Morgan", "−0.041", "中段 CI 多不含 0"],
            ["DRFP", "−0.036", "与 Morgan 高度同向"],
            ["DFT", "−0.036", "表示间相关 r≳0.86"],
        ],
    )
    add_p(doc, "B=30 同向略负；B=100 压缩至约 −0.01（冷启动接近天花板）。", indent=False)

    add_h(doc, "4.2 S0：匹配目标板初始化", 2)
    add_p(doc, "目的：排除「冷/池化 init 不匹配」造成的伪负迁移。", indent=False)
    add_fig(
        doc,
        FIGS / "fig_edbo_suzuki_s0_vs_main_pair_delta.png",
        "图3　S0 匹配 init 稳健性",
    )
    add_bul(doc, "设计：先采目标板 init；源采样 RNG 与目标分离 → 同 (target, seed) 匹配 100%。")
    add_bul(doc, "B=40：Morgan Main −0.041 → S0 −0.034；DFT Main −0.036 → S0 −0.030。")
    add_bul(doc, "分析：略减负，符号不翻；init 错配会放大，但不是负迁移主因。")

    add_h(doc, "4.3 A2：任务内秩 / 百分位池化", 2)
    add_p(
        doc,
        "目的：若伤害主要来自不同任务产率标度不可比，则秩变换应明显修复。"
        "设计：任务内将 y 转为百分位秩，再在秩空间做与 A1 相同的无 task ID 池化。",
    )
    add_bul(doc, "B=40 vs S0 cold：A1 raw Morgan/DFT ≈ −0.034/−0.030；A2 rank ≈ −0.026/−0.029。")
    add_bul(doc, "略减负、未翻正 → 排除「只是产量标度」的单一解释。")

    add_h(doc, "4.4 A3：源标签噪声降权", 2)
    add_p(doc, "目的：中等源降权能否在无 task ID 设定下稳定扭转负迁移？", indent=False)
    add_fig(
        doc,
        FIGS / "fig_edbo_suzuki_ladder_A1A2A3_B40.png",
        "图4　A1–A3 方法阶梯（B=40，相对 S0 cold）",
    )
    add_bul(doc, "设计：α_src=1e-4/w_s，w∈{0.1, 0.25, 0.5}。")
    add_bul(doc, "B=40：A3 各档 ≈ A1（约 −0.03x），几乎重合。")
    add_bul(doc, "口径：测试档未修复；不是「降权永远无效」。")

    add_h(doc, "五、数据分析", 1)
    add_h(doc, "5.1 为何主窗口是 B=30–50", 2)
    add_fig(
        doc,
        ESI / "fig_edbo_suzuki_headroom_vs_delta_frac_B40.png",
        "图5　headroom 与 Δfrac（B=40）",
    )
    add_bul(doc, "多表示同向 → 不是单一指纹特例。")
    add_bul(doc, "中段伤害最大，此时仍有 headroom；B=100 终值差异被天花板压缩。")
    add_bul(doc, "故「接近零@100」≠「池化无害」。")

    add_h(doc, "5.2 A3 健全性", 2)
    add_bul(doc, "实现核对：权重进入 GP 对角噪声，接线真实，非空转。")
    add_bul(doc, "后验探针：w=0.1 vs 1 时 max|Δμ| 约 0.05；极端小 w 才明显靠近 cold。")
    add_bul(doc, "轨迹：约 73%–78% 的 pair×seed 在 w=0.1–0.5 间采集路径相同。")
    add_p(
        doc,
        "结论口径：中等噪声膨胀未能一致优于原始池化；早期 n_s≫n_t 时源点仍被高度信任。"
        "不把 w→0（≈退回冷启动）当作有效方法臂。",
    )

    add_h(doc, "5.3 Pair 异质性", 2)
    add_fig(
        doc,
        ESI / "fig_edbo_suzuki_morgan_pair_delta_frac_heatmap_B40.png",
        "图6　pair 层面 Δfrac 热图（示例表示，B=40）",
    )
    add_bul(doc, "效应高度 pair 依赖；正尾薄、负尾厚。")
    add_bul(doc, "不是「一个坏 pair」，也不是「平均正、偶发负」。")
    add_bul(doc, "实践上不能指望「换个源就默认变好」。")

    add_h(doc, "5.4 假设检验综合", 2)
    add_tbl(
        doc,
        ["对照", "检验假设", "结果", "排除"],
        [
            ["A0 vs A1", "源/目标 y 可交换", "中段平均负迁移", "原始标签可交换"],
            ["S0", "伤害仅来自 init 错配", "仍负，略减轻", "init 为唯一原因"],
            ["A1 vs A2", "主要是产率标度", "略减负，未翻正", "纯标度故事"],
            ["A1 vs A3", "中等可靠性收缩够用", "≈A1，未稳定修复", "简单降权即够"],
        ],
    )
    add_p(
        doc,
        "综合解读（非唯一证明）：与 task mismatch 一致——共享 X ⇏ 跨任务历史 y 可交换。",
        bold=True,
        indent=False,
    )

    add_h(doc, "六、边界对照：Doyle2018 · PK2022 · amination", 1)
    add_bul(doc, "Doyle2018（Science 2018 BH CN）：共享 L×B×Add；label(OHE) pair Δ≈+0.064 → 外部边界，别处可正。")
    add_bul(doc, "PK2022（Science 2022；CHAOS 四板）：一维添加剂库；label pair Δ≈+0.13~+0.15 → 1D 边界，不可当多因子主证据。")
    add_bul(doc, "xTB 全量 / 边界库多表示重跑：本阶段不做。")

    add_h(doc, "6.1 EDBO amination min S0（全量 2560）", 2)
    add_p(
        doc,
        "协议对齐 Suzuki S0：cold vs label_warm；Morgan + DFT；8 底物；20 seeds；"
        "主终点为 B∈{30,40,50} 的 mean Δfrac。",
        indent=False,
    )
    add_tbl(
        doc,
        ["表示", "mean Δfrac", "median", "n_pos / n_neg / n_near0（56 pairs）"],
        [
            ["DFT", "+0.043", "+0.029", "33 / 8 / 15"],
            ["Morgan", "+0.013", "+0.005", "15 / 13 / 28"],
            ["Suzuki C1（对照）", "≈ −0.034 ~ −0.037", "—", "正尾薄、负尾厚"],
        ],
    )
    add_bul(doc, "DFT 平均略正；Morgan 一半 near0，谈不上可靠增益。")
    add_bul(
        doc,
        "异质性：正均值被 target=sub_s4 拉高（DFT≈+0.20）。去掉 s4 后 DFT≈+0.020、Morgan≈0。",
    )
    add_bul(
        doc,
        "写作口径：不做跨家族「一致不安全默认」升级；amination 证明效应可 family-dependent；"
        "主主张仍锚定多因子 Suzuki。",
    )

    add_h(doc, "七、阶段结论与解读边界", 1)
    add_p(
        doc,
        "在主验证多因子共享库上：无 task ID 的持续历史标签池化，相对冷启动不是安全默认。"
        "多表示同向负；匹配 init / 秩 / 名义降权未翻正。"
        "第二家族 amination 平均偏正/近零且高度靶依赖 → 跨家族不可一刀切。",
    )
    add_p(doc, "当前支持的解读：", bold=True, indent=False)
    add_bul(doc, "共享 X ≠ 历史 y 可默认混用。")
    add_bul(doc, "与 task mismatch 相一致。")
    add_bul(doc, "别处/他家族可正 → 非普适禁令。")
    add_p(doc, "不宜外推：", bold=True, indent=False)
    add_bul(doc, "历史数据没用 / 一切 transfer 无效。")
    add_bul(doc, "task-aware / MTGP 也必然失败。")
    add_bul(doc, "任意反应家族都必然负迁移。")

    add_h(doc, "八、进度与待讨论", 1)
    add_tbl(
        doc,
        ["事项", "状态"],
        [
            ["C1 / S0 / A2 / A3 全套验证", "完成"],
            ["Task-level 重推断 + A2/A3 审计", "完成"],
            ["相似度机制（top-k ↔ Δfrac）", "完成"],
            ["Amination 全量 2560", "完成；family-dependent"],
            ["S5 / task-ID 策略 + 小试点", "完成（未扩全网格）"],
            ["xTB 全量 / 边界多表示重跑", "本阶段不做"],
        ],
    )
    add_p(doc, "待讨论：", bold=True, indent=False)
    add_bul(doc, "Amination 作为边界写入正文/SI，是否同意不做跨家族强升级？")
    add_bul(doc, "是否扩跑 S5（仅 init）全网格，作为「然后怎么办」主对照？")
    add_bul(doc, "轻量 task-ID 试点≈A1；是否上更强 MTGP，还是先写清局限？")
    add_bul(doc, "边界库是否需要多表示重跑，还是维持现状作对照即可？")

    add_h(doc, "九、材料路径", 1)
    add_bul(doc, "本汇报 PPT/DOCX：docs/briefings/（v0.7）")
    add_bul(doc, "数字汇总：results/paper_stats/EXPERIMENT_SUMMARY.md")
    add_bul(doc, "Amination：results/paper_stats/edbo_amination_min_s0_SUMMARY.md")
    add_bul(doc, "结果目录：results/external_* · results/transfer_grid*")
    add_bul(doc, "图件：docs/figs/main/ 与 docs/figs/")

    add_p(doc, "—— 汇报完 ——", indent=False, center=True)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    doc.save(DOC_OUT)
    try:
        import shutil

        shutil.copy2(DOC_OUT, DOC_ALIAS)
    except Exception:
        pass
    return DOC_OUT


def main() -> int:
    import sys

    only = {a.lower() for a in sys.argv[1:]}
    if "--docx-only" in only or "docx" in only:
        doc = build_docx()
        print("DOCX", doc)
        return 0
    ppt = build_ppt()
    doc = build_docx()
    print("PPT ", ppt)
    print("DOCX", doc)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
