# -*- coding: utf-8 -*-
"""生成 BOUSE 介绍 PPT（明亮风格，16:9，原生形状可编辑）。

用法::

    python slides/build_slides.py

输出::

    slides/BOUSE_intro.pptx
    slides/assets/opt_progress.png
"""
from __future__ import annotations

import math
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent
BOUSE = ROOT.parent
ASSETS = ROOT / "assets"
ASSETS.mkdir(exist_ok=True)
OUT = ROOT / "BOUSE_intro.pptx"

# ---------------------------------------------------------------------------
# 调色板（明亮）
# ---------------------------------------------------------------------------
INK = RGBColor(0x1E, 0x29, 0x3B)      # slate-800
SUB = RGBColor(0x64, 0x74, 0x8B)      # slate-500
FAINT = RGBColor(0x94, 0xA3, 0xB8)    # slate-400
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

BLUE = RGBColor(0x25, 0x63, 0xEB)
SKY = RGBColor(0x0E, 0xA5, 0xE9)
TEAL = RGBColor(0x14, 0xB8, 0xA6)
GREEN = RGBColor(0x10, 0xB9, 0x81)
ORANGE = RGBColor(0xF9, 0x73, 0x16)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
RED = RGBColor(0xEF, 0x44, 0x44)

BLUE_L = RGBColor(0xEF, 0xF6, 0xFF)
SKY_L = RGBColor(0xF0, 0xF9, 0xFF)
TEAL_L = RGBColor(0xF0, 0xFD, 0xFA)
GREEN_L = RGBColor(0xEC, 0xFD, 0xF5)
ORANGE_L = RGBColor(0xFF, 0xF7, 0xED)
AMBER_L = RGBColor(0xFF, 0xFB, 0xEB)
PURPLE_L = RGBColor(0xF5, 0xF3, 0xFF)
RED_L = RGBColor(0xFE, 0xF2, 0xF2)
SLATE_L = RGBColor(0xF8, 0xFA, 0xFC)

BLUE_B = RGBColor(0xBF, 0xDB, 0xFE)
SKY_B = RGBColor(0xBA, 0xE6, 0xFD)
TEAL_B = RGBColor(0x99, 0xF6, 0xE4)
GREEN_B = RGBColor(0xA7, 0xF3, 0xD0)
ORANGE_B = RGBColor(0xFE, 0xD7, 0xAA)
AMBER_B = RGBColor(0xFD, 0xE6, 0x8A)
PURPLE_B = RGBColor(0xDD, 0xD6, 0xFE)
SLATE_B = RGBColor(0xE2, 0xE8, 0xF0)

RING = RGBColor(0xDB, 0xEA, 0xFE)
GRID_DOT = RGBColor(0xCB, 0xD5, 0xE1)

FONT = "Microsoft YaHei"
MONO = "Consolas"

PAGE_W, PAGE_H = 13.333, 7.5

# ---------------------------------------------------------------------------
# 基础助手
# ---------------------------------------------------------------------------

def _noshadow(shape):
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def box(slide, x, y, w, h, fill=None, line=None, line_w=1.0,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=None, dash=None):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    _noshadow(sp)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
        if dash:
            sp.line.dash_style = dash
    sp.text_frame.word_wrap = True
    return sp


def txt(slide, x, y, w, h, text, size=12, color=INK, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, font=FONT, line_spacing=1.0, space_after=0):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        if space_after:
            p.space_after = Pt(space_after)
        r = p.add_run()
        r.text = ln
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        # 中文字体（east-asian typeface）
        rPr = r.font._rPr
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            ea = rPr.makeelement(qn("a:ea"), {})
            rPr.append(ea)
        ea.set("typeface", font)
    return tb


def shape_text(sp, text, size=12, color=INK, bold=False, align=PP_ALIGN.CENTER,
               anchor=MSO_ANCHOR.MIDDLE, font=FONT, line_spacing=1.0):
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Pt(2) if "left" in m or "right" in m else Pt(1))
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = ln
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        rPr = r.font._rPr
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            ea = rPr.makeelement(qn("a:ea"), {})
            rPr.append(ea)
        ea.set("typeface", font)


def chip(slide, x, y, w, h, text, fill, fg, size=10.5, bold=False, line=None,
         radius=0.5, font=FONT):
    sp = box(slide, x, y, w, h, fill=fill, line=line, radius=radius)
    shape_text(sp, text, size=size, color=fg, bold=bold, font=font)
    return sp


def conn(slide, x1, y1, x2, y2, color=FAINT, w=1.5, dash=None, arrow=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(w)
    if dash:
        c.line.dash_style = dash
    _noshadow(c)
    if arrow:
        ln = c.line._get_or_add_ln()
        el = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
        ln.append(el)
    return c


def dot(slide, cx, cy, d, fill, line=None):
    return box(slide, cx - d / 2, cy - d / 2, d, d, fill=fill, line=line,
               shape=MSO_SHAPE.OVAL)


def header(slide, kicker, title, accent, kicker_fill):
    box(slide, 0.55, 0.5, 0.09, 0.66, fill=accent, shape=MSO_SHAPE.RECTANGLE)
    chip(slide, 0.78, 0.5, 1.5, 0.3, kicker, kicker_fill, accent, size=11, bold=True)
    txt(slide, 0.78, 0.84, 11.0, 0.5, title, size=25, color=INK, bold=True)


def footer(slide, idx):
    txt(slide, 0.55, 7.08, 4.0, 0.3, "BOUSE · AI4CHEM", size=9, color=FAINT)
    txt(slide, 12.3, 7.08, 0.5, 0.3, str(idx), size=9, color=FAINT, align=PP_ALIGN.RIGHT)


def file_row(slide, x, y, w, name, note, accent, h=0.52):
    box(slide, x, y, w, h, fill=WHITE, line=SLATE_B, radius=0.18)
    ic = box(slide, x + 0.12, y + 0.09, 0.34, 0.34, fill=accent, shape=MSO_SHAPE.FOLDED_CORNER)
    txt(slide, x + 0.55, y + 0.05, 2.15, 0.42, name, size=9.5, color=INK, bold=True,
        font=MONO, anchor=MSO_ANCHOR.MIDDLE)
    txt(slide, x + 2.72, y + 0.05, w - 2.85, 0.42, note, size=8.5, color=SUB,
        anchor=MSO_ANCHOR.MIDDLE)


# ---------------------------------------------------------------------------
# 图表 PNG（真实 demo 数据）
# ---------------------------------------------------------------------------

def render_chart() -> Path:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    plt.rcParams["font.sans-serif"] = ["Microsoft YaHei", "SimHei"]
    plt.rcParams["axes.unicode_minus"] = False

    csv = BOUSE / "edbo" / "workspaces" / "deoxy_demo" / "history.csv"
    df = pd.read_csv(csv)
    y = df["yield"].tolist()
    x = list(range(1, len(y) + 1))
    best, run = [], -10**9
    for v in y:
        run = max(run, v)
        best.append(run)

    fig, ax = plt.subplots(figsize=(4.35, 3.0), dpi=220)
    ax.step(x, best, where="post", color="#2563EB", lw=2, label="历史最优")
    ax.scatter(x, y, s=64, color="#F97316", zorder=3, label="单次实验")
    ax.annotate(f"{best[-1]}%", (x[-1], best[-1]), xytext=(7, -3),
                textcoords="offset points", color="#1E293B",
                fontsize=11, fontweight="bold")
    ax.set_xlabel("实验序号", fontsize=10.5, color="#475569")
    ax.set_ylabel("收率 / %", fontsize=10.5, color="#475569")
    ax.set_xticks(x)
    ax.set_ylim(0, max(best) * 1.35 + 4)
    ax.grid(axis="y", color="#E5E7EB", lw=0.8)
    ax.set_axisbelow(True)
    for s in ("top", "right", "left"):
        ax.spines[s].set_visible(False)
    ax.spines["bottom"].set_color("#CBD5E1")
    ax.tick_params(colors="#64748B", labelsize=9.5)
    ax.legend(loc="upper left", frameon=False, fontsize=9.5)
    fig.tight_layout(pad=0.6)
    out = ASSETS / "opt_progress.png"
    fig.savefig(out, facecolor="white")
    plt.close(fig)
    return out


# ---------------------------------------------------------------------------
# S1 封面
# ---------------------------------------------------------------------------

def s1(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # 背景装饰
    dot(s, 11.6, 0.9, 3.4, SKY_L)
    dot(s, 12.6, 5.9, 2.2, ORANGE_L)
    dot(s, 0.4, 6.6, 2.6, BLUE_L)
    ring = box(s, 9.0, 4.6, 2.4, 2.4, fill=None, line=BLUE_B, line_w=2, shape=MSO_SHAPE.OVAL)
    dot(s, 10.2, 4.35, 0.28, BLUE)
    dot(s, 11.35, 5.3, 0.28, ORANGE)
    dot(s, 10.9, 6.9, 0.28, GREEN)
    dot(s, 9.15, 6.2, 0.28, PURPLE)

    chip(s, 0.9, 1.55, 2.9, 0.42, "AI4CHEM · 反应优化工作区", BLUE_L, BLUE, size=12, bold=True)
    txt(s, 0.88, 2.15, 8.5, 1.6, "BOUSE", size=80, color=INK, bold=True)
    box(s, 0.95, 3.62, 1.15, 0.1, fill=ORANGE, shape=MSO_SHAPE.RECTANGLE)
    box(s, 2.2, 3.62, 0.35, 0.1, fill=BLUE, shape=MSO_SHAPE.RECTANGLE)
    txt(s, 0.92, 3.9, 9.0, 0.6, "闭环反应优化：描述符生成 ∥ EDBO 贝叶斯优化", size=21, color=SUB)

    labels = [("① 描述符生成", BLUE_L, BLUE), ("② EDBO 优化向导", ORANGE_L, ORANGE),
              ("③ 契约文件对接", GREEN_L, GREEN)]
    x = 0.92
    for t, f, c in labels:
        chip(s, x, 4.85, 2.35, 0.5, t, f, c, size=13, bold=True)
        x += 2.55
    txt(s, 0.92, 6.5, 8.0, 0.4, "汇报人：zhangzhou · 2026-07-24", size=12, color=FAINT)


# ---------------------------------------------------------------------------
# S2 总体架构
# ---------------------------------------------------------------------------

def s2(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "架构", "双界面并行 · 只通过约定文件对接", BLUE, BLUE_L)

    # 左：描述符界面
    box(s, 0.55, 1.7, 3.7, 3.9, fill=BLUE_L, line=BLUE_B, radius=0.06)
    chip(s, 0.85, 1.95, 2.2, 0.42, "① 描述符生成", WHITE, BLUE, size=13, bold=True)
    txt(s, 0.85, 2.42, 3.2, 0.3, "Streamlit 界面 · :8502", size=9.5, color=SUB)
    chip(s, 0.85, 2.86, 3.1, 0.42, "输入：SMILES 列表", WHITE, INK, size=10.5)
    for i, b in enumerate(["RDKit 2D", "MACCS", "Morgan", "Mordred"]):
        chip(s, 0.85 + (i % 2) * 1.6, 3.44 + (i // 2) * 0.52, 1.5, 0.4, b,
             WHITE, SKY, size=10, bold=True, line=SKY_B)
    chip(s, 0.85, 4.62, 3.1, 0.42, "输出：descriptor_<因子>.csv", WHITE, INK, size=10.5)
    txt(s, 0.85, 5.12, 3.2, 0.3, "另有 CLI：python cli.py", size=9, color=SUB)

    # 右：EDBO 向导
    box(s, 9.05, 1.7, 3.7, 3.9, fill=ORANGE_L, line=ORANGE_B, radius=0.06)
    chip(s, 9.35, 1.95, 2.4, 0.42, "② EDBO 优化向导", WHITE, ORANGE, size=13, bold=True)
    txt(s, 9.35, 2.42, 3.2, 0.3, "Streamlit 界面 · :8501", size=9.5, color=SUB)
    steps = ["1  项目与目标", "2  定义搜索域", "3  推荐下一轮", "4  回填结果"]
    for i, t in enumerate(steps):
        chip(s, 9.35, 2.86 + i * 0.56, 3.1, 0.44, t, WHITE, INK, size=10.5)
        if i < 3:
            txt(s, 9.35 + 1.45, 3.28 + i * 0.56, 0.3, 0.14, "↓", size=9, color=FAINT)

    # 中间：契约文件
    txt(s, 4.62, 1.85, 4.1, 0.3, "约定文件交接", size=12, color=SUB, align=PP_ALIGN.CENTER, bold=True)
    mid = box(s, 5.0, 2.3, 3.35, 2.2, fill=WHITE, line=GREEN_B, radius=0.08)
    box(s, 5.0, 2.3, 3.35, 0.14, fill=GREEN, shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE, radius=0.5)
    for i, (n, c) in enumerate([("CONTRACT.md", GREEN), ("descriptor_<因子>.csv", BLUE),
                                ("handoff.py 校验/导入", PURPLE)]):
        chip(s, 5.25, 2.62 + i * 0.6, 2.85, 0.46, n, SLATE_L, INK, size=10, font=MONO)
    conn(s, 4.25, 3.4, 5.0, 3.4, color=GREEN, w=2.5, arrow=True)
    conn(s, 8.35, 3.4, 9.05, 3.4, color=GREEN, w=2.5, arrow=True)

    # 底部：共享工作区
    box(s, 0.55, 5.85, 12.2, 0.95, fill=SLATE_L, line=SLATE_B, radius=0.12)
    box(s, 0.85, 6.08, 0.55, 0.42, fill=AMBER, shape=MSO_SHAPE.FOLDED_CORNER)
    txt(s, 1.6, 6.02, 10.9, 0.35, "共享工作区  edbo/workspaces/<项目>/", size=12.5, color=INK, bold=True)
    txt(s, 1.6, 6.38, 10.9, 0.3, "config.json · descriptor_*.csv · history.csv · last_recommendations.csv",
        size=10.5, color=SUB, font=MONO)
    footer(s, 2)


# ---------------------------------------------------------------------------
# S3 闭环工作流
# ---------------------------------------------------------------------------

def s3(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "流程", "闭环工作流：推荐 → 实验 → 回填 → 再推荐", ORANGE, ORANGE_L)

    cx, cy, r = 6.667, 4.32, 1.95
    # 环
    ring = box(s, cx - r, cy - r, 2 * r, 2 * r, fill=RING, shape=MSO_SHAPE.DONUT)
    try:
        ring.adjustments[0] = 0.045
    except Exception:
        pass
    # 方向三角
    for i in range(6):
        mid = math.radians(-90 + i * 60 + 30)
        px, py = cx + r * math.cos(mid), cy + r * math.sin(mid)
        tri = box(s, px - 0.11, py - 0.11, 0.22, 0.22, fill=BLUE,
                  shape=MSO_SHAPE.ISOSCELES_TRIANGLE)
        dx, dy = -math.sin(mid), math.cos(mid)
        tri.rotation = math.degrees(math.atan2(dx, -dy))

    nodes = [("准备\nSMILES", BLUE), ("生成\n描述符", SKY), ("对接\n导入域", PURPLE),
             ("BO 推荐\n下一轮", ORANGE), ("湿实验\n测定", GREEN), ("回填\n历史", AMBER)]
    for i, (label, color) in enumerate(nodes):
        th = math.radians(-90 + i * 60)
        nx, ny = cx + r * math.cos(th), cy + r * math.sin(th)
        c = dot(s, nx, ny, 1.14, color)
        shape_text(c, label, size=11.5, color=WHITE, bold=True, line_spacing=0.95)
        badge = dot(s, nx - 0.44, ny - 0.44, 0.34, WHITE, line=color)
        shape_text(badge, str(i + 1), size=11, color=color, bold=True)

    # 中心
    c0 = dot(s, cx, cy, 1.75, WHITE, line=BLUE_B)
    shape_text(c0, "闭环优化\n目标：收率↑", size=13.5, color=INK, bold=True, line_spacing=1.1)

    chip(s, 0.85, 6.35, 2.6, 0.44, "①–③  一次性准备", BLUE_L, BLUE, size=11, bold=True)
    chip(s, 9.9, 6.35, 2.6, 0.44, "④–⑥  每轮迭代", ORANGE_L, ORANGE, size=11, bold=True)
    chip(s, 9.9, 1.7, 2.6, 0.44, "⑤ 人在回路 · 真实实验", GREEN_L, GREEN, size=10.5, bold=True)
    footer(s, 3)


# ---------------------------------------------------------------------------
# S4 描述符生成
# ---------------------------------------------------------------------------

def s4(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "模块 ①", "描述符生成：SMILES → 数值特征表", SKY, SKY_L)

    # 输入
    box(s, 0.55, 2.6, 2.3, 1.9, fill=WHITE, line=SLATE_B, radius=0.08)
    chip(s, 0.8, 2.82, 1.8, 0.4, "molecules.csv", SLATE_L, INK, size=10.5, bold=True, font=MONO)
    txt(s, 0.8, 3.32, 1.9, 0.9, "smiles\nCCO\nc1ccccc1\nO=S(F)=O…", size=9.5,
        color=SUB, font=MONO, line_spacing=1.25)
    txt(s, 0.55, 2.14, 2.3, 0.3, "输入", size=11, color=SUB, bold=True, align=PP_ALIGN.CENTER)

    # 后端容器
    box(s, 3.75, 1.95, 4.4, 4.55, fill=SLATE_L, line=SLATE_B, radius=0.05)
    chip(s, 4.0, 2.12, 1.7, 0.38, "generators/", WHITE, SUB, size=10.5, bold=True, font=MONO)
    backs = [("RDKit 2D", "~200 维理化描述符", BLUE, BLUE_L),
             ("MACCS", "167 bit 结构密钥", SKY, SKY_L),
             ("Morgan", "ECFP 指纹 · 位数可配", TEAL, TEAL_L),
             ("Mordred", "1800+ 描述符（可选）", PURPLE, PURPLE_L),
             ("Clean", "外部表清洗 · 降维（DFT）", AMBER, AMBER_L)]
    for i, (name, note, c, cl) in enumerate(backs):
        y = 2.66 + i * 0.76
        box(s, 4.0, y, 3.9, 0.62, fill=WHITE, line=SLATE_B, radius=0.16)
        dot(s, 4.32, y + 0.31, 0.2, c)
        txt(s, 4.52, y + 0.06, 1.35, 0.5, name, size=11.5, color=INK, bold=True,
            anchor=MSO_ANCHOR.MIDDLE)
        txt(s, 5.85, y + 0.06, 2.0, 0.5, note, size=9, color=SUB, anchor=MSO_ANCHOR.MIDDLE)
    conn(s, 2.85, 3.55, 3.75, 3.55, color=FAINT, w=2, arrow=True)

    # 输出
    txt(s, 9.0, 2.14, 3.6, 0.3, "输出（契约格式）", size=11, color=SUB, bold=True,
        align=PP_ALIGN.CENTER)
    box(s, 9.0, 2.6, 3.8, 2.35, fill=BLUE_L, line=BLUE_B, radius=0.08)
    chip(s, 9.25, 2.78, 3.3, 0.42, "descriptor_solvent.csv", WHITE, BLUE, size=11,
         bold=True, font=MONO)
    tbl_shape = s.shapes.add_table(3, 4, Inches(9.25), Inches(3.36), Inches(3.3), Inches(1.35))
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(1.3)
    for j in range(1, 4):
        tbl.columns[j].width = Inches(0.66)
    hdr = ["molecule_id", "f1", "f2", "…"]
    rows = [["C1CCOC1", "0.31", "1.07", "…"], ["O=CN(C)C", "0.12", "0.88", "…"]]
    for j, t in enumerate(hdr):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        cell.text = t
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.size = Pt(9)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = WHITE
        p.runs[0].font.name = MONO
    for i, row in enumerate(rows):
        for j, t in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            cell.text = t
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.runs[0].font.size = Pt(9)
            p.runs[0].font.color.rgb = SUB
            p.runs[0].font.name = MONO
    chip(s, 9.25, 5.18, 3.3, 0.42, "失败分子 → *_failed.csv（不进 EDBO）", RED_L, RED, size=9.5)
    conn(s, 8.15, 3.55, 9.0, 3.55, color=FAINT, w=2, arrow=True)

    chip(s, 0.55, 6.75, 3.6, 0.42, "入口 A：Streamlit 界面 :8502", SKY_L, SKY, size=10.5, bold=True)
    chip(s, 4.35, 6.75, 3.6, 0.42, "入口 B：python cli.py（同一后端）", SLATE_L, SUB,
         size=10.5, bold=True, font=MONO)
    footer(s, 4)


# ---------------------------------------------------------------------------
# S5 EDBO 向导四步
# ---------------------------------------------------------------------------

def s5(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "模块 ②", "EDBO 优化向导：四步完成一轮", ORANGE, ORANGE_L)

    cards = [
        ("项目与目标", ORANGE, ORANGE_L, ORANGE_B,
         ["目标列：yield（最大化）", "batch = 5 / 轮", "内置模板一键建档"]),
        ("定义搜索域", BLUE, BLUE_L, BLUE_B,
         ["化学因子 → 描述符表", "数值因子 → 取值列表", "全空间 31 万 → 采样 2500"]),
        ("推荐下一轮", PURPLE, PURPLE_L, PURPLE_B,
         ["EI / TS / UCB / PI", "无模型：LHS · Sobol", "随机 · 极大极小距离"]),
        ("回填结果", GREEN, GREEN_L, GREEN_B,
         ["导出测定模板 CSV", "填入实验收率", "合并进 history.csv"]),
    ]
    x0, y0, w, h, gap = 0.55, 1.95, 2.83, 3.9, 0.3
    for i, (name, c, cl, cb, lines) in enumerate(cards):
        x = x0 + i * (w + gap)
        box(s, x, y0, w, h, fill=cl, line=cb, radius=0.07)
        num = dot(s, x + w / 2, y0 + 0.62, 0.85, c)
        shape_text(num, str(i + 1), size=26, color=WHITE, bold=True)
        txt(s, x + 0.2, y0 + 1.2, w - 0.4, 0.45, name, size=16, color=INK, bold=True,
            align=PP_ALIGN.CENTER)
        for j, ln in enumerate(lines):
            box(s, x + 0.25, y0 + 1.85 + j * 0.58, w - 0.5, 0.46, fill=WHITE, radius=0.22)
            txt(s, x + 0.4, y0 + 1.85 + j * 0.58, w - 0.7, 0.46, ln, size=10, color=SUB,
                anchor=MSO_ANCHOR.MIDDLE)
        if i < 3:
            ch = box(s, x + w - 0.03, y0 + 1.7, 0.38, 0.5, fill=FAINT, shape=MSO_SHAPE.CHEVRON)

    # 回路
    y_back = 6.35
    conn(s, x0 + 3 * (w + gap) + w / 2, y0 + h, x0 + 3 * (w + gap) + w / 2, y_back,
         color=GREEN, w=2)
    conn(s, x0 + 3 * (w + gap) + w / 2, y_back, x0 + 2 * (w + gap) + w / 2, y_back,
         color=GREEN, w=2)
    conn(s, x0 + 2 * (w + gap) + w / 2, y_back, x0 + 2 * (w + gap) + w / 2, y0 + h + 0.02,
         color=GREEN, w=2, arrow=True)
    chip(s, 5.75, y_back - 0.21, 1.9, 0.42, "回填后再推荐", GREEN_L, GREEN, size=10.5, bold=True)
    footer(s, 5)


# ---------------------------------------------------------------------------
# S6 Demo 案例
# ---------------------------------------------------------------------------

def s6(prs, chart: Path):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "案例", "Demo：脱氧氟化反应优化（EDBO 官方数据）", GREEN, GREEN_L)

    # 反应式
    box(s, 0.55, 1.8, 7.5, 1.75, fill=SLATE_L, line=SLATE_B, radius=0.08)
    chip(s, 0.8, 2.0, 1.6, 0.55, "R–OH\n醇", WHITE, INK, size=11, bold=True)
    txt(s, 2.42, 2.0, 0.35, 0.55, "+", size=20, color=SUB, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    chip(s, 2.8, 2.0, 1.85, 0.55, "Ar–SO₂F\n磺酰氟", WHITE, INK, size=11, bold=True)
    conn(s, 4.85, 2.28, 6.05, 2.28, color=INK, w=2.25, arrow=True)
    chip(s, 6.2, 2.0, 1.6, 0.55, "R–F\n氟化物", GREEN_L, GREEN, size=11, bold=True)
    txt(s, 0.8, 3.02, 7.0, 0.35, "Deoxyfluorination · 碱 / 溶剂 / 温度 / 浓度 / 当量",
        size=9.5, color=SUB)

    # 因子
    box(s, 0.55, 3.8, 7.5, 1.5, fill=WHITE, line=SLATE_B, radius=0.08)
    txt(s, 0.8, 3.92, 3.0, 0.3, "7 个因子", size=11.5, color=INK, bold=True)
    chem = [("磺酰氟 ×10", BLUE, BLUE_L), ("碱 ×10", SKY, SKY_L), ("溶剂 ×5", TEAL, TEAL_L)]
    for i, (t, c, cl) in enumerate(chem):
        chip(s, 0.8 + i * 1.52, 4.26, 1.42, 0.38, t, cl, c, size=10, bold=True)
    txt(s, 5.5, 4.26, 2.4, 0.38, "化学因子 · 描述符 20 维", size=8.5, color=FAINT,
        anchor=MSO_ANCHOR.MIDDLE)
    nums = ["浓度 ×5", "当量 ×5 ×2", "温度 ×5"]
    for i, t in enumerate(nums):
        chip(s, 0.8 + i * 1.52, 4.76, 1.42, 0.38, t, SLATE_L, SUB, size=10, bold=True)
    txt(s, 5.5, 4.76, 2.4, 0.38, "数值因子 · 各 5 档", size=8.5, color=FAINT,
        anchor=MSO_ANCHOR.MIDDLE)

    # 搜索空间
    box(s, 0.55, 5.5, 7.5, 1.35, fill=WHITE, line=SLATE_B, radius=0.08)
    txt(s, 0.8, 5.62, 3.0, 0.3, "搜索空间", size=11.5, color=INK, bold=True)
    for r in range(3):
        for c in range(14):
            dot(s, 1.0 + c * 0.17, 6.12 + r * 0.17, 0.09, GRID_DOT)
    chip(s, 0.8, 6.5, 2.55, 0.3, "全空间 312,500 组", SLATE_L, SUB, size=9)
    conn(s, 3.6, 6.32, 4.35, 6.32, color=FAINT, w=2, arrow=True)
    for r in range(3):
        for c in range(8):
            dot(s, 4.6 + c * 0.17, 6.12 + r * 0.17, 0.09, ORANGE)
    chip(s, 4.45, 6.5, 1.6, 0.3, "采样 2,500", ORANGE_L, ORANGE, size=9, bold=True)
    txt(s, 6.15, 6.5, 1.8, 0.3, "保留全部历史点", size=8.5, color=FAINT)

    # 图表
    box(s, 8.3, 1.8, 4.45, 4.05, fill=WHITE, line=SLATE_B, radius=0.06)
    txt(s, 8.55, 1.98, 4.0, 0.32, "优化进程（demo 前 5 组实验）", size=11.5, color=INK, bold=True)
    s.shapes.add_picture(str(chart), Inches(8.42), Inches(2.35), width=Inches(4.2))

    # 指标条
    stats = [("batch", "5"), ("采集函数", "EI"), ("训练", "100 iters"), ("目标", "yield 最大化")]
    for i, (k, v) in enumerate(stats):
        x = 8.3 + i * 1.16
        box(s, x, 6.05, 1.06, 0.8, fill=GREEN_L, radius=0.14)
        txt(s, x, 6.14, 1.06, 0.3, k, size=8.5, color=GREEN, align=PP_ALIGN.CENTER, bold=True)
        txt(s, x, 6.4, 1.06, 0.36, v, size=10.5, color=INK, align=PP_ALIGN.CENTER, bold=True)
    footer(s, 6)


# ---------------------------------------------------------------------------
# S7 文件契约
# ---------------------------------------------------------------------------

def s7(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "契约", "交接契约：只交换文件，不共享 UI", PURPLE, PURPLE_L)

    # 左：规则
    txt(s, 0.55, 1.8, 3.4, 0.35, "描述符 CSV 规则", size=13, color=INK, bold=True)
    rules = ["必需列 molecule_id（唯一）", "特征列全部为数值", "文件名 descriptor_<因子>.csv",
             "因子 key 与 config.json 一致", "UTF-8 · 首行表头"]
    for i, t in enumerate(rules):
        box(s, 0.55, 2.3 + i * 0.62, 3.5, 0.5, fill=WHITE, line=SLATE_B, radius=0.2)
        dot(s, 0.82, 2.55 + i * 0.62, 0.14, PURPLE)
        txt(s, 1.0, 2.3 + i * 0.62, 3.0, 0.5, t, size=10.5, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # 中：工作区文件树
    box(s, 4.45, 1.8, 4.35, 4.9, fill=SLATE_L, line=SLATE_B, radius=0.06)
    box(s, 4.75, 2.05, 0.5, 0.38, fill=AMBER, shape=MSO_SHAPE.FOLDED_CORNER)
    txt(s, 5.4, 2.05, 3.2, 0.38, "edbo/workspaces/deoxy_demo/", size=11, color=INK,
        bold=True, font=MONO, anchor=MSO_ANCHOR.MIDDLE)
    files = [("config.json", "因子·目标·配置", BLUE),
             ("descriptor_base.csv", "特征表 20 维", SKY),
             ("descriptor_solvent.csv", "按因子一份", SKY),
             ("history.csv", "实验 + yield", GREEN),
             ("last_recommendations.csv", "最近建议", ORANGE)]
    for i, (n, note, c) in enumerate(files):
        file_row(s, 4.75, 2.62 + i * 0.62, 3.75, n, note, c)
    chip(s, 4.75, 5.9, 3.75, 0.5, "levels_<因子>.csv（OHE 编码时可选）", WHITE, SUB, size=9.5,
         line=SLATE_B)

    # 右：校验/导入
    txt(s, 9.2, 1.8, 3.6, 0.35, "校验与导入", size=13, color=INK, bold=True)
    box(s, 9.2, 2.3, 3.6, 2.3, fill=WHITE, line=SLATE_B, radius=0.08)
    txt(s, 9.45, 2.5, 3.2, 1.9,
        "python scripts/\\\n  validate_handoff.py \\\n  --workspace …/deoxy_demo\n\n"
        "python scripts/\\\n  import_descriptor.py …",
        size=9.5, color=SUB, font=MONO, line_spacing=1.3)
    chip(s, 9.2, 4.85, 3.6, 0.48, "UI 内也可一键导入（推荐）", PURPLE_L, PURPLE, size=10.5, bold=True)
    box(s, 9.2, 5.55, 3.6, 1.15, fill=PURPLE_L, line=PURPLE_B, radius=0.1)
    txt(s, 9.45, 5.72, 3.15, 0.85, "契约化 = 可插拔\n任何 BO 策略都能接入同一工作区",
        size=11, color=INK, bold=True, line_spacing=1.25)
    footer(s, 7)


# ---------------------------------------------------------------------------
# S8 启动与进展
# ---------------------------------------------------------------------------

def s8(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "使用", "一键启动 · 进展与路线", AMBER, AMBER_L)

    # 左：启动
    box(s, 0.55, 1.95, 5.3, 1.15, fill=AMBER_L, line=AMBER_B, radius=0.1)
    dot(s, 1.15, 2.52, 0.55, AMBER)
    txt(s, 0.93, 2.36, 0.45, 0.22, "▶", size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, 1.55, 2.18, 4.1, 0.35, "start_bouse.bat", size=15, color=INK, bold=True, font=MONO)
    txt(s, 1.55, 2.58, 4.1, 0.3, "双击 · 同时打开两个界面", size=10, color=SUB)

    def browser(y, port, name, c, cl):
        box(s, 0.55, y, 5.3, 1.25, fill=WHITE, line=SLATE_B, radius=0.09)
        box(s, 0.55, y, 5.3, 0.34, fill=SLATE_L, shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE,
            radius=0.28)
        for k, dc in enumerate([RED, AMBER, GREEN]):
            dot(s, 0.82 + k * 0.22, y + 0.17, 0.11, dc)
        chip(s, 2.0, y + 0.05, 2.6, 0.24, f"localhost:{port}", WHITE, SUB, size=8.5, font=MONO,
             line=SLATE_B)
        chip(s, 0.85, y + 0.52, 2.4, 0.5, name, cl, c, size=12.5, bold=True)
        txt(s, 3.45, y + 0.52, 2.2, 0.5, "Streamlit", size=9.5, color=FAINT,
            anchor=MSO_ANCHOR.MIDDLE)

    conn(s, 3.2, 3.1, 3.2, 3.42, color=FAINT, w=2, arrow=True)
    browser(3.45, "8501", "② EDBO 优化向导", ORANGE, ORANGE_L)
    browser(4.95, "8502", "① 描述符生成", BLUE, BLUE_L)

    # 右：进展
    txt(s, 6.45, 1.95, 3.0, 0.35, "已完成", size=13, color=GREEN, bold=True)
    done = ["EDBO 可视化向导（四步）", "描述符生成界面 + CLI", "契约对接 · 校验 · 导入"]
    for i, t in enumerate(done):
        box(s, 6.45, 2.4 + i * 0.6, 6.3, 0.48, fill=GREEN_L, radius=0.2)
        chk = dot(s, 6.78, 2.64 + i * 0.6, 0.3, GREEN)
        shape_text(chk, "✓", size=12, color=WHITE, bold=True)
        txt(s, 7.05, 2.4 + i * 0.6, 5.5, 0.48, t, size=11.5, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    txt(s, 6.45, 4.45, 3.0, 0.35, "下一步", size=13, color=AMBER, bold=True)
    nxt = ["更多描述符后端", "DFT 表 → clean 管线", "接入更多 BO 策略", "自动化闭环"]
    for i, t in enumerate(nxt):
        chip(s, 6.45 + (i % 2) * 3.25, 4.9 + (i // 2) * 0.6, 3.1, 0.48, t, WHITE, SUB,
             size=11, line=AMBER_B)

    box(s, 0.55, 6.38, 12.2, 0.58, fill=INK, radius=0.16)
    txt(s, 0.55, 6.38, 12.2, 0.58, "BOUSE — 让反应优化的闭环转起来", size=15, color=WHITE,
        bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 8)


# ---------------------------------------------------------------------------

def main():
    chart = render_chart()
    prs = Presentation()
    prs.slide_width = Inches(PAGE_W)
    prs.slide_height = Inches(PAGE_H)
    s1(prs)
    s2(prs)
    s3(prs)
    s4(prs)
    s5(prs)
    s6(prs, chart)
    s7(prs)
    s8(prs)
    prs.save(OUT)
    print(f"saved: {OUT} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
